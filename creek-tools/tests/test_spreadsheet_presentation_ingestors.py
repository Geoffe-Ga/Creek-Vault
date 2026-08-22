"""Tests for the spreadsheet and presentation ingestors (#57)."""

from __future__ import annotations

from typing import TYPE_CHECKING, Any

import frontmatter
import pytest

from creek.ingest.base import assemble_ingested_fragment
from creek.ingest.encoding import UndecodableBytesError
from creek.ingest.presentations import (
    PRESENTATION_EXTENSIONS,
    PresentationBackend,
    PresentationData,
    PresentationIngestor,
    PythonPptxBackend,
    PythonPptxUnavailableError,
    SlideData,
)
from creek.ingest.spreadsheets import (
    SPREADSHEET_EXTENSIONS,
    OpenpyxlBackend,
    OpenpyxlUnavailableError,
    SheetData,
    SpreadsheetBackend,
    SpreadsheetIngestor,
    WorkbookData,
)
from creek.models import SourcePlatform
from creek.vault.writer import VaultWriter

if TYPE_CHECKING:
    from pathlib import Path

    from creek.ingest.base import IngestedFragment, Ingestor, ParsedFragment


# ---- Stub backends -----------------------------------------------------


class StubSpreadsheetBackend:
    """Deterministic spreadsheet backend for tests."""

    def __init__(self, workbooks: dict[str, WorkbookData] | None = None) -> None:
        """Seed the backend with canned workbook data keyed by filename."""
        self._workbooks = dict(workbooks or {})
        self.calls: list[str] = []

    def is_available(self) -> bool:
        """Stubs are always available."""
        return True

    def read_workbook(
        self,
        path: Path,
        *,
        has_header: bool | None = None,
    ) -> WorkbookData:
        """Return the canned workbook for *path*.

        ``has_header`` is accepted for Protocol compatibility (#165)
        but ignored — canned ``WorkbookData`` already has its
        ``SheetData.headers`` baked in by each test.
        """
        del has_header
        self.calls.append(f"read_workbook:{path.name}")
        return self._workbooks[path.name]


class StubPresentationBackend:
    """Deterministic presentation backend for tests."""

    def __init__(
        self,
        presentations: dict[str, PresentationData] | None = None,
    ) -> None:
        """Seed the backend with canned presentation data keyed by filename."""
        self._presentations = dict(presentations or {})
        self.calls: list[str] = []

    def is_available(self) -> bool:
        """Stubs are always available."""
        return True

    def read_presentation(self, path: Path) -> PresentationData:
        """Return the canned presentation for *path*."""
        self.calls.append(f"read_presentation:{path.name}")
        return self._presentations[path.name]


def _make_import_blocker(blocked: set[str]) -> object:
    """Build a ``__import__`` replacement that raises for *blocked*."""
    import builtins

    real_import = builtins.__import__

    def _blocked_import(
        name: str,
        module_globals: dict[str, object] | None = None,
        module_locals: dict[str, object] | None = None,
        fromlist: tuple[str, ...] = (),
        level: int = 0,
    ) -> object:
        """Raise ImportError for blocked roots; defer everything else."""
        root = name.split(".", 1)[0]
        if root in blocked or name in blocked:
            msg = f"mocked-missing: {name}"
            raise ImportError(msg)
        return real_import(name, module_globals, module_locals, fromlist, level)

    return _blocked_import


# ---- Helpers -----------------------------------------------------------


def _write_xlsx_placeholder(path: Path) -> None:
    """Write a minimal XLSX magic-byte placeholder (parsed via stub backend)."""
    path.write_bytes(b"PK\x03\x04 placeholder")


def _write_csv(path: Path, rows: list[list[str]]) -> None:
    """Write a comma-separated value file at *path*."""
    import csv

    with path.open("w", encoding="utf-8", newline="") as handle:
        csv.writer(handle).writerows(rows)


def _write_pptx_placeholder(path: Path) -> None:
    """Write a minimal PPTX magic-byte placeholder (parsed via stub backend)."""
    path.write_bytes(b"PK\x03\x04 placeholder pptx")


# =======================================================================
# SpreadsheetIngestor
# =======================================================================


class TestModuleConstants:
    """Public constants exposed by the spreadsheet module."""

    def test_extensions_cover_xlsx_and_csv(self) -> None:
        """Both .xlsx and .csv are recognised."""
        assert ".xlsx" in SPREADSHEET_EXTENSIONS
        assert ".csv" in SPREADSHEET_EXTENSIONS


class TestSheetData:
    """Invariants on the :class:`SheetData` value object."""

    def test_immutable_rows_are_tuples(self) -> None:
        """Rows are stored as tuples of tuples — immutable end-to-end."""
        sheet = SheetData(
            name="Q1",
            headers=("month", "amount"),
            rows=(("Jan", "100"), ("Feb", "200")),
        )
        assert isinstance(sheet.rows, tuple)
        assert isinstance(sheet.rows[0], tuple)


# ---- discover ---------------------------------------------------------


class TestSpreadsheetDiscover:
    """`SpreadsheetIngestor.discover` finds the right files."""

    def test_finds_xlsx_and_csv(self, tmp_path: Path) -> None:
        """Both .xlsx and .csv files are discovered."""
        _write_xlsx_placeholder(tmp_path / "a.xlsx")
        _write_csv(tmp_path / "b.csv", [["h"], ["v"]])
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        raws = ingestor.discover(tmp_path)
        assert {raw.path.name for raw in raws} == {"a.xlsx", "b.csv"}

    def test_ignores_non_spreadsheet_files(self, tmp_path: Path) -> None:
        """Files with unsupported extensions are skipped."""
        _write_xlsx_placeholder(tmp_path / "ok.xlsx")
        (tmp_path / "note.md").write_text("# md")
        (tmp_path / "shot.png").write_bytes(b"\x89PNG")
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        raws = ingestor.discover(tmp_path)
        assert [raw.path.name for raw in raws] == ["ok.xlsx"]

    def test_recursive_discovery(self, tmp_path: Path) -> None:
        """Files in subdirectories are found."""
        nested = tmp_path / "deep"
        nested.mkdir()
        _write_xlsx_placeholder(nested / "deep.xlsx")
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        raws = ingestor.discover(tmp_path)
        assert any("deep.xlsx" in str(raw.path) for raw in raws)

    def test_extension_match_is_case_insensitive(self, tmp_path: Path) -> None:
        """``.XLSX`` and ``.CSV`` count too."""
        _write_xlsx_placeholder(tmp_path / "shout.XLSX")
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        raws = ingestor.discover(tmp_path)
        assert {raw.path.name for raw in raws} == {"shout.XLSX"}

    def test_single_image_file_returns_empty(self, tmp_path: Path) -> None:
        """A single non-spreadsheet file returns an empty discovery."""
        png = tmp_path / "x.png"
        png.write_bytes(b"\x89PNG")
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        assert ingestor.discover(png) == []

    def test_does_not_load_file_bytes(self, tmp_path: Path) -> None:
        """Discovery records the path but does not slurp the bytes."""
        path = tmp_path / "huge.xlsx"
        _write_xlsx_placeholder(path)
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        raws = ingestor.discover(tmp_path)
        assert raws[0].content == b""


# ---- parse / ingest ---------------------------------------------------


class TestSpreadsheetParse:
    """Each sheet becomes a :class:`ParsedFragment`."""

    def test_one_fragment_per_sheet(self, tmp_path: Path) -> None:
        """A workbook with two sheets yields two fragments."""
        path = tmp_path / "report.xlsx"
        _write_xlsx_placeholder(path)
        backend = StubSpreadsheetBackend(
            workbooks={
                "report.xlsx": WorkbookData(
                    sheets=(
                        SheetData(
                            name="Q1",
                            headers=("month", "amount"),
                            rows=(("Jan", "100"),),
                        ),
                        SheetData(
                            name="Q2",
                            headers=("month", "amount"),
                            rows=(("Apr", "300"),),
                        ),
                    ),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        raws = ingestor.discover(tmp_path)
        fragments = ingestor.parse(raws[0])
        assert len(fragments) == 2
        assert {f.metadata["sheet"] for f in fragments} == {"Q1", "Q2"}

    def test_fragment_metadata_records_dimensions(self, tmp_path: Path) -> None:
        """``rows`` / ``columns`` are recorded on each fragment's metadata."""
        path = tmp_path / "data.xlsx"
        _write_xlsx_placeholder(path)
        backend = StubSpreadsheetBackend(
            workbooks={
                "data.xlsx": WorkbookData(
                    sheets=(
                        SheetData(
                            name="Sheet1",
                            headers=("a", "b", "c"),
                            rows=(("1", "2", "3"), ("4", "5", "6")),
                        ),
                    ),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert fragments[0].metadata["rows"] == 2
        assert fragments[0].metadata["columns"] == 3

    def test_empty_sheet_yields_no_fragment(self, tmp_path: Path) -> None:
        """A sheet with no headers and no rows is dropped."""
        path = tmp_path / "blank.xlsx"
        _write_xlsx_placeholder(path)
        backend = StubSpreadsheetBackend(
            workbooks={
                "blank.xlsx": WorkbookData(
                    sheets=(SheetData(name="Empty", headers=None, rows=()),),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert fragments == []

    def test_parsed_fragment_carries_typed_sheetdata_payload(
        self,
        tmp_path: Path,
    ) -> None:
        """Sheet structure is carried as a typed :class:`SheetData`, not a dict.

        Issue #166 — the pre-refactor code flattened the backend's
        ``SheetData`` into ``metadata["headers"]`` /
        ``metadata["row_data"]`` and ``convert_to_markdown`` re-read
        them. The new contract puts the typed object on
        :attr:`ParsedFragment.payload` so the schema lives in one
        place and mypy can catch a rename at the call site.
        """
        path = tmp_path / "typed.xlsx"
        _write_xlsx_placeholder(path)
        original = SheetData(
            name="Typed",
            headers=("a", "b"),
            rows=(("1", "2"),),
        )
        backend = StubSpreadsheetBackend(
            workbooks={"typed.xlsx": WorkbookData(sheets=(original,))},
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert fragments[0].payload is original
        # The structured fields are NOT mirrored back into the dict,
        # so a refactor that drops the mirror later does not regress
        # this test.
        assert "headers" not in fragments[0].metadata
        assert "row_data" not in fragments[0].metadata


# ---- convert_to_markdown ----------------------------------------------


class TestSpreadsheetMarkdown:
    """Markdown rendering uses GFM tables."""

    def test_small_sheet_renders_full_table(self, tmp_path: Path) -> None:
        """A small sheet (≤ summary threshold) renders every row."""
        path = tmp_path / "small.xlsx"
        _write_xlsx_placeholder(path)
        backend = StubSpreadsheetBackend(
            workbooks={
                "small.xlsx": WorkbookData(
                    sheets=(
                        SheetData(
                            name="Sheet1",
                            headers=("month", "amount"),
                            rows=(("Jan", "100"), ("Feb", "200")),
                        ),
                    ),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "| month | amount |" in markdown
        assert "| --- | --- |" in markdown
        assert "| Jan | 100 |" in markdown
        assert "| Feb | 200 |" in markdown

    def test_large_sheet_is_summarised(self, tmp_path: Path) -> None:
        """A sheet > summary threshold shows first/last N rows + total count."""
        path = tmp_path / "huge.xlsx"
        _write_xlsx_placeholder(path)
        big_rows = tuple((str(i), f"v{i}") for i in range(150))
        backend = StubSpreadsheetBackend(
            workbooks={
                "huge.xlsx": WorkbookData(
                    sheets=(
                        SheetData(
                            name="Big",
                            headers=("idx", "value"),
                            rows=big_rows,
                        ),
                    ),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "150 rows" in markdown  # Total count surfaced.
        assert "v0" in markdown  # First few rows kept.
        assert "v149" in markdown  # Last few rows kept.
        # Middle rows trimmed:
        assert "v75" not in markdown

    def test_large_sheet_renders_single_header_after_summary(
        self,
        tmp_path: Path,
    ) -> None:
        """Summary note precedes a single GFM table — no duplicate header."""
        path = tmp_path / "huge.xlsx"
        _write_xlsx_placeholder(path)
        big_rows = tuple((str(i), f"v{i}") for i in range(150))
        backend = StubSpreadsheetBackend(
            workbooks={
                "huge.xlsx": WorkbookData(
                    sheets=(
                        SheetData(
                            name="Big",
                            headers=("idx", "value"),
                            rows=big_rows,
                        ),
                    ),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        header_line = "| idx | value |"
        # The header row should appear exactly once.
        assert markdown.count(header_line) == 1
        # The separator row should appear exactly once too.
        assert markdown.count("| --- | --- |") == 1
        # The summary note should appear before the header in the document.
        summary_idx = markdown.index("Showing first")
        header_idx = markdown.index(header_line)
        assert summary_idx < header_idx, (
            "Summary note should precede the (single) header row."
        )

    def test_pipe_and_newline_in_cell_are_escaped(self, tmp_path: Path) -> None:
        """Cells containing ``|`` or newlines do not corrupt GFM tables."""
        path = tmp_path / "escapes.xlsx"
        _write_xlsx_placeholder(path)
        backend = StubSpreadsheetBackend(
            workbooks={
                "escapes.xlsx": WorkbookData(
                    sheets=(
                        SheetData(
                            name="Sheet1",
                            headers=("a|b", "c"),
                            rows=(("x | y", "line1\nline2"),),
                        ),
                    ),
                ),
            },
        )
        ingestor = SpreadsheetIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        # The literal cell text must be present in escaped form.
        assert r"a\|b" in markdown
        assert r"x \| y" in markdown
        # Newlines inside a cell must not split the table row.
        assert "line1<br>line2" in markdown
        # Every non-empty markdown line that begins with `|` should have
        # the same column count (3 pipes for 2 columns).
        table_rows = [line for line in markdown.splitlines() if line.startswith("|")]
        column_counts = {line.count("|") - line.count(r"\|") for line in table_rows}
        assert column_counts == {3}, (
            f"Mismatched pipe counts indicate broken escaping: {column_counts}"
        )


# ---- generate_frontmatter ---------------------------------------------


class TestSpreadsheetFrontmatter:
    """Frontmatter shape for spreadsheet fragments."""

    def test_frontmatter_records_platform_and_dimensions(
        self,
        tmp_path: Path,
    ) -> None:
        """The ingestor's RETURNED dict carries platform plus sheet/row/col counts.

        Scope is the point of this docstring. This asserts the mapping
        ``generate_frontmatter`` returns — no ``VaultWriter``, no file, no
        disk. Until #1392 it read like coverage of the vault's frontmatter
        and was not: all three keys were dropped on the way to disk by
        ``Fragment``'s ``extra="ignore"``, and nothing here would have
        noticed. The disk half is
        ``TestMultiSheetWorkbookIdentity.test_sheet_rows_and_columns_survive_to_disk``;
        read the two together or neither means what it appears to.
        """
        from datetime import UTC, datetime

        from creek.ingest.base import ParsedFragment

        fragment = ParsedFragment(
            content="",
            metadata={
                "original_file": str(tmp_path / "x.xlsx"),
                "sheet": "Q1",
                "rows": 5,
                "columns": 3,
            },
            source_path=str(tmp_path / "x.xlsx"),
            timestamp=datetime(2026, 4, 27, tzinfo=UTC),
        )
        ingestor = SpreadsheetIngestor(backend=StubSpreadsheetBackend())
        fm = ingestor.generate_frontmatter(fragment)
        assert fm["source"]["platform"] == SourcePlatform.SPREADSHEET.value
        assert fm["sheet"] == "Q1"
        assert fm["rows"] == 5
        assert fm["columns"] == 3


class TestSpreadsheetAuthoredAt:
    """FEAT-031: XLSX core properties ``created`` / ``modified`` → ``authored_at``."""

    def _write_real_xlsx(
        self,
        path: Path,
        *,
        created: object | None = None,
        modified: object | None = None,
    ) -> None:
        """Write a real, openpyxl-parseable XLSX with the given core properties.

        The factory uses the genuine openpyxl writer rather than a
        placeholder so :func:`_extract_xlsx_authored_at` reads the
        round-tripped core properties as it would in production.
        """
        import openpyxl

        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "Sheet1"
        worksheet.append(["h1", "h2"])
        worksheet.append(["a", "b"])
        if created is not None:
            workbook.properties.created = created  # type: ignore[assignment]
        if modified is not None:
            workbook.properties.modified = modified  # type: ignore[assignment]
        workbook.save(path)

    def test_authored_at_from_xlsx_created(self, tmp_path: Path) -> None:
        """The ``created`` core property is promoted to ``authored_at``."""
        from datetime import UTC, datetime

        xlsx = tmp_path / "report.xlsx"
        self._write_real_xlsx(
            xlsx,
            created=datetime(2024, 3, 15, 8, 30, tzinfo=UTC),
        )
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(xlsx)[0])
        assert len(fragments) == 1
        authored = fragments[0].metadata["authored_at"]
        assert authored is not None
        assert authored.date() == datetime(2024, 3, 15).date()
        assert authored.tzinfo is not None

    def test_authored_at_falls_through_to_modified(self, tmp_path: Path) -> None:
        """No ``created`` → use ``modified`` as the next-best source date."""
        from datetime import UTC, datetime

        xlsx = tmp_path / "report.xlsx"
        # openpyxl always emits a default ``created`` if we don't override
        # it, so set both explicitly: created=None, modified=our date.
        self._write_real_xlsx(
            xlsx,
            created=None,
            modified=datetime(2024, 4, 1, tzinfo=UTC),
        )
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(xlsx)[0])
        authored = fragments[0].metadata["authored_at"]
        # openpyxl may auto-fill ``created`` even when set to ``None``;
        # accept either ``modified`` or whatever the library produces,
        # as long as it's tz-aware and not silently dropped.
        assert authored is not None
        assert authored.tzinfo is not None

    def test_csv_has_no_authored_at(self, tmp_path: Path) -> None:
        """CSV files carry no core properties → ``authored_at`` is ``None``."""
        csv_path = tmp_path / "rows.csv"
        _write_csv(csv_path, [["a", "b"], ["1", "2"]])
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(csv_path)[0])
        assert len(fragments) == 1
        assert fragments[0].metadata["authored_at"] is None

    def test_authored_at_in_frontmatter_when_present(self, tmp_path: Path) -> None:
        """Generated frontmatter carries ``authored_at`` as ISO string."""
        from datetime import UTC, datetime

        xlsx = tmp_path / "report.xlsx"
        self._write_real_xlsx(
            xlsx,
            created=datetime(2024, 3, 15, tzinfo=UTC),
        )
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(xlsx)[0])
        fm = ingestor.generate_frontmatter(fragments[0])
        assert "authored_at" in fm
        assert fm["authored_at"].startswith("2024-03-15")

    def test_no_authored_at_omits_key_from_frontmatter(self, tmp_path: Path) -> None:
        """When ``authored_at`` is ``None`` the key is absent (terse YAML)."""
        csv_path = tmp_path / "rows.csv"
        _write_csv(csv_path, [["a"], ["1"]])
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(csv_path)[0])
        fm = ingestor.generate_frontmatter(fragments[0])
        assert "authored_at" not in fm


# ---- CSV path --------------------------------------------------------


class TestCsvFiles:
    """CSV files are read directly through Python's stdlib."""

    def test_csv_renders_as_single_sheet(self, tmp_path: Path) -> None:
        """A CSV file is parsed as one sheet named after the file."""
        csv_path = tmp_path / "rows.csv"
        _write_csv(csv_path, [["a", "b"], ["1", "2"], ["3", "4"]])
        # Use the production OpenpyxlBackend; CSV parsing does not
        # require the optional openpyxl dep.
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert len(fragments) == 1
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "| a | b |" in markdown
        assert "| 1 | 2 |" in markdown

    def test_csv_with_utf8_bom_is_decoded(self, tmp_path: Path) -> None:
        """CSV files saved with a UTF-8 BOM (common from Excel) decode cleanly."""
        csv_path = tmp_path / "bom.csv"
        # Real-world Excel "CSV UTF-8" exports prepend the BOM.
        csv_path.write_bytes(
            "﻿name,city\nAlice,Münster\nBob,São Paulo\n".encode(),
        )
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        # Header row must not carry the BOM character.
        assert "﻿" not in markdown
        assert "| name | city |" in markdown
        assert "Münster" in markdown
        assert "São Paulo" in markdown

    def test_csv_with_cp1252_falls_back(self, tmp_path: Path) -> None:
        """CSV files saved as CP1252 (legacy Excel default) decode without crashing."""
        csv_path = tmp_path / "legacy.csv"
        csv_path.write_bytes("name,note\nAlice,café\n".encode("cp1252"))
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "| name | note |" in markdown
        assert "café" in markdown

    def test_csv_shift_jis_decodes_with_chardet(self, tmp_path: Path) -> None:
        """A Shift-JIS CSV is decoded via the chardet probe (BUG-010).

        Without the probe, ``cp1252`` would silently turn the multi-byte
        Japanese sequences into mojibake. ``chardet`` needs a generous
        sample to reach the confidence threshold, so this test uses
        a long, varied corpus rather than a couple of rows.
        """
        csv_path = tmp_path / "shift_jis.csv"
        # A varied phrase repeated until chardet has hundreds of bytes
        # of distinguishing signal — small samples produce low
        # confidence and would (correctly) fall through to cp1252.
        phrase = (
            "名前,都市,メモ\n"
            "田中,東京,ありがとうございます\n"
            "鈴木,大阪,初めまして、よろしくお願いします\n"
            "佐藤,京都,日本語のテキストはマルチバイトです\n"
        )
        rows = phrase * 60
        csv_path.write_bytes(rows.encode("shift_jis"))
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        # If we had decoded as cp1252 the kanji would be replaced by
        # Latin-1 characters. The exact byte-for-byte output depends on
        # chardet's choice (Shift-JIS / cp932 are equivalent for these
        # characters), so the test asserts that at least one of the
        # input kanji round-trips intact.
        assert any(ch in markdown for ch in "東京名前田中大阪京都")

    def test_csv_shift_jis_cells_all_survive_the_detection_tables(
        self,
        tmp_path: Path,
    ) -> None:
        """Every Shift-JIS cell reaches the markdown intact (#1257).

        chardet 7.4.3 -> 7.6.0 reshapes the statistical tables behind
        :func:`chardet.detect`, and this is one of the two call sites
        where that choice decides what text enters a vault. Shift-JIS is
        the family that clears
        :data:`CSV_CHARDET_CONFIDENCE_THRESHOLD` (CP932 @ 0.93 on this
        corpus, on both versions), so it is the one that actually
        exercises the detection branch rather than the cp1252 fallback.

        The neighbouring ``test_csv_shift_jis_decodes_with_chardet``
        asserts that *at least one* kanji survives, which a partially
        wrong codec still passes. This asserts every cell, because a
        vault of half-corrupted rows looks healthy until someone reads
        it.

        The codec *name* is deliberately not asserted: 7.4.3 answers
        ``cp932`` and 7.6.0 answers ``CP932``. Both decode identically
        and ``_read_csv`` lowercases before its reject-set test, so
        pinning the name would fail on a harmless rename while still
        passing on real mojibake.
        """
        csv_path = tmp_path / "shift_jis_cells.csv"
        rows = (
            "名前,都市,メモ\n"
            "田中,東京,ありがとうございます\n"
            "鈴木,大阪,初めまして、よろしくお願いします\n"
            "佐藤,京都,日本語のテキストはマルチバイトです\n"
        )
        csv_path.write_bytes((rows * 60).encode("shift_jis"))
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        for cell in (
            "名前",
            "都市",
            "メモ",
            "田中",
            "東京",
            "ありがとうございます",
            "初めまして、よろしくお願いします",
            "日本語のテキストはマルチバイトです",
        ):
            assert cell in markdown, (
                f"Shift-JIS cell {cell!r} did not survive decoding; "
                "chardet chose a codec whose tables disagree with the "
                "writer's, which reaches a vault as mojibake (#1257)"
            )

    def test_csv_cp1252_fallback_logs_warning(
        self,
        tmp_path: Path,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        """The cp1252 fallback logs a WARNING with the file path (BUG-010).

        Pure-ASCII content fits every codec, so chardet declines to
        commit (its top guess is ``ascii``, which we reject in favour
        of cp1252 to keep behaviour stable). The warning lets the user
        spot mojibake before it lands in the vault.
        """
        import logging

        csv_path = tmp_path / "ascii.csv"
        # Bytes invalid as utf-8 force the fallback path even on
        # shorter samples where chardet is uncertain.
        csv_path.write_bytes(b"name,note\r\nAlice,caf\xe9\r\n")
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        with caplog.at_level(logging.WARNING, logger="creek.ingest.spreadsheets"):
            ingestor.parse(ingestor.discover(tmp_path)[0])
        assert any(
            "decoded as cp1252" in record.message and "ascii.csv" in record.message
            for record in caplog.records
        )


# ---- the CSV encoding decision matrix (#1589, #1591) ------------------

#: One row per encoding family that reaches ``_read_csv``, with the
#: chardet verdict measured on *this* corpus under chardet 7.6.0. The
#: figures are recorded here so nobody has to re-derive them, and
#: because both issues were filed with wrong ones: #1589 claimed
#: GB18030 scored 0.54-0.63 and that CSV commas diluted it, and that
#: Shift-JIS was unaffected. Measured, none of that holds — Chinese
#: prose scores *lower* than the same characters in a CSV, confidence
#: is flat with length, and Japanese lands under the gate too.
#:
#: What actually splits the rows is two independent facts: whether the
#: corpus scores under :data:`CSV_CHARDET_CONFIDENCE_THRESHOLD`, and
#: whether its bytes happen to include one of cp1252's five undefined
#: values (0x81 0x8d 0x8f 0x90 0x9d). The first face is silent
#: mojibake (#1589); the second is a ``UnicodeDecodeError`` that
#: ``_parse_safe`` turns into a dropped file (#1591). Same gate, two
#: symptoms, so they are one table rather than two tests.
#:
#:   id                  codec      chardet verdict     was
#:   zh-gbk              gbk        GB18030 @ 0.38      mojibake
#:   zh-rare-lead-bytes  gbk        GB18030 @ 0.57      crash on 0x81
#:   jp-short-sjis       shift_jis  CP932   @ 0.32      crash on 0x81
#:   jp-long-sjis        shift_jis  CP932   @ 0.89      correct
#:   ko-euckr           euc-kr     CP949   @ 0.69      mojibake
#:   cp1252-rich         cp1252     cp862   @ 0.08      correct
#:   utf8-plain          utf-8      never reaches chardet
#:   utf8-bom            utf-8-sig  never reaches chardet
_CSV_ENCODING_MATRIX: tuple[tuple[str, str, str], ...] = (
    (
        "zh-gbk",
        "gbk",
        "姓名,城市,备注\n"
        "王伟,北京,他每天早上都会沿着小河散步\n"
        "李娜,上海,河水清澈岸边的树木在风中轻轻摇动\n",
    ),
    (
        # Contains 丂丄丅丆丏 — GBK encodes these with a 0x81 lead
        # byte, which cp1252 leaves undefined. This row is #1591's
        # crash wearing Chinese rather than Japanese, which is why
        # "the defect is specific to the Chinese detector" and "a
        # Japanese spreadsheet aborts the run" are both wrong about
        # which axis matters.
        "zh-rare-lead-bytes",
        "gbk",
        "姓名,城市,备注\n"
        "王伟,北京,他每天早上都会沿着小河散步丂丄丅丆丏\n"
        "李娜,上海,河水清澈岸边的树木在风中轻轻摇动\n",
    ),
    (
        # 、 is 0x8141 and 。 is 0x8142 in Shift-JIS: the exact bytes
        # #1591 was filed on. Two short rows, so it scores far under
        # the gate — character variety, not length, is what moves
        # chardet's confidence.
        "jp-short-sjis",
        "shift_jis",
        "名前、年齢。\n田中、三十。\n",
    ),
    (
        # The one CJK corpus that already clears the gate, kept as a
        # regression guard: the fix must not disturb the path that
        # was already correct.
        "jp-long-sjis",
        "shift_jis",
        "名前,都市,メモ\n"
        "田中,東京,ありがとうございます\n"
        "鈴木,大阪,初めまして、よろしくお願いします\n",
    ),
    (
        "ko-euckr",
        "euc-kr",
        "이름,도시,비고\n"
        "김철수,서울,그는 매일 아침 강변을 걷습니다\n"
        "이영희,부산,그녀는 숲을 좋아합니다\n",
    ),
    (
        # The row that forbids the two fixes both issues suggested.
        # chardet answers a *single-byte* codec here (cp862 @ 0.08),
        # and raw.decode() of that guess succeeds while producing
        # naïve -> naïve-shaped garbage. So "trust a detection that
        # round-trips, even below the gate" (#1589) would corrupt
        # this row, and swapping the fallback to latin-1 (#1591)
        # would turn every smart quote, dash and € into a control
        # character. Both stay green only because the fix keys on
        # codec *class*, not on confidence or decodability.
        "cp1252-rich",
        "cp1252",
        "name,note\nAlice,café naïve — “quoted”\nBob,résumé \u2013 £85 €9 ©\n",
    ),
    ("utf8-plain", "utf-8", "name,city\n田中,東京\nAlice,Münster\n"),
    ("utf8-bom", "utf-8-sig", "name,city\nMünster,São Paulo\n"),
)


def _cells(text: str) -> list[str]:
    """Split a CSV corpus into its individual cell strings."""
    return [
        cell for line in text.strip().split("\n") for cell in line.split(",") if cell
    ]


class TestCsvEncodingDecisionMatrix:
    """Every encoding family round-trips through ``_read_csv`` intact.

    The pre-existing CSV encoding tests assert substrings — "café" is
    in the markdown, or *at least one* kanji survived. Both survive the
    corruption they exist to catch, because a wrong codec still leaves
    the ASCII run and some fraction of the multi-byte run intact. These
    assert **every cell**, which is the only assertion a half-corrupted
    vault fails.
    """

    @pytest.mark.parametrize(
        ("codec", "corpus"),
        [(codec, corpus) for _, codec, corpus in _CSV_ENCODING_MATRIX],
        ids=[name for name, _, _ in _CSV_ENCODING_MATRIX],
    )
    def test_every_cell_survives_decoding(
        self,
        tmp_path: Path,
        codec: str,
        corpus: str,
    ) -> None:
        """Bytes written in *codec* reach the markdown unchanged.

        Args:
            tmp_path: Per-test directory holding the CSV.
            codec: Codec the corpus is written in.
            corpus: Source CSV text.
        """
        csv_path = tmp_path / "matrix.csv"
        csv_path.write_bytes(corpus.encode(codec))
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "\ufeff" not in markdown, (
            "a BOM reached the rendered table; utf-8-sig must strip it"
        )
        for cell in _cells(corpus):
            assert cell in markdown, (
                f"cell {cell!r} did not survive decoding of {codec} bytes. "
                "Either chardet's guess was rejected and cp1252 mangled "
                "the row (#1589), or the codec chosen disagrees with the "
                "writer's. A vault of half-corrupted rows looks healthy "
                "until someone reads it."
            )

    def test_the_matrix_still_covers_every_family(self) -> None:
        """The table has not been emptied, trimmed, or re-scoped.

        Deleting a row removes a guard without turning anything red —
        the parametrisation simply runs fewer cases and the gate stays
        green. The two rows that matter most are the ones a naive fix
        would drop: ``cp1252-rich`` is what forbids the latin-1 swap
        and the round-trip tie-breaker, and ``jp-long-sjis`` is the
        only corpus that was already correct.
        """
        assert {name for name, _, _ in _CSV_ENCODING_MATRIX} == {
            "zh-gbk",
            "zh-rare-lead-bytes",
            "jp-short-sjis",
            "jp-long-sjis",
            "ko-euckr",
            "cp1252-rich",
            "utf8-plain",
            "utf8-bom",
        }

    def test_undecodable_binary_csv_fails_loudly_and_writes_nothing(
        self,
        tmp_path: Path,
    ) -> None:
        """A binary ``.csv`` raises, and ``ingest()`` drops it with an error.

        #1591 suggested falling back to latin-1, which decodes all 256
        byte values — that would convert this loud failure into a
        silent fragment of garbage. The fallback chain therefore
        refuses binary input before it reaches any all-accepting
        codec, so a genuinely undecodable file stays loud.
        """
        csv_path = tmp_path / "binary.csv"
        csv_path.write_bytes(bytes(range(256)) * 4)
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        with pytest.raises(UndecodableBytesError):
            ingestor.parse(ingestor.discover(tmp_path)[0])
        result = ingestor.ingest(tmp_path)
        assert result.fragments == [], (
            "a binary CSV produced a fragment; undecodable input must "
            "never reach the vault"
        )
        assert any("binary.csv" in error for error in result.errors), (
            "the dropped file was not recorded in result.errors, so the "
            "operator has no signal that a file went missing (#1591)"
        )


# ---- has_header override (#165) ---------------------------------------


class TestHasHeaderOverride:
    """Caller-controllable header detection (#165).

    The legacy heuristic — "first row qualifies as headers when every
    cell is a non-empty string" — produces false positives on data
    sheets whose first row happens to be all-strings (country names,
    labels, etc.). Three modes:

    * ``has_header=None`` (default) preserves the heuristic.
    * ``has_header=True``  forces the first row to be headers.
    * ``has_header=False`` forces the first row to be data and the
      sheet's headers to be auto-generated ``colN`` placeholders at
      render time.
    """

    def test_csv_heuristic_promotes_all_string_first_row(
        self,
        tmp_path: Path,
    ) -> None:
        """Default ``has_header=None`` still auto-detects (heuristic intact).

        This pins today's behaviour: a CSV whose first row is all
        non-empty strings is treated as a header row. The next test
        shows why this is sometimes wrong; together they describe the
        bug the override fixes.
        """
        csv_path = tmp_path / "countries.csv"
        _write_csv(
            csv_path, [["France", "Germany", "Spain"], ["Italy", "Greece", "Portugal"]]
        )
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        # The heuristic promotes the first row, so "France" lands in headers.
        sheet = fragments[0].payload
        assert isinstance(sheet, SheetData)
        assert list(sheet.headers or ()) == ["France", "Germany", "Spain"]
        assert [list(row) for row in sheet.rows] == [
            ["Italy", "Greece", "Portugal"],
        ]

    def test_csv_has_header_false_keeps_first_row_as_data(
        self,
        tmp_path: Path,
    ) -> None:
        """``has_header=False`` overrides the heuristic for the false-positive case.

        Same all-string-first-row sheet as above, but now the caller
        knows it's data — the first row stays in ``row_data`` and the
        rendered table uses generated ``col1..colN`` headers.
        """
        csv_path = tmp_path / "countries.csv"
        _write_csv(
            csv_path, [["France", "Germany", "Spain"], ["Italy", "Greece", "Portugal"]]
        )
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        raws = ingestor.discover(tmp_path)
        fragments = ingestor.parse(raws[0], has_header=False)
        sheet = fragments[0].payload
        assert isinstance(sheet, SheetData)
        assert sheet.headers in (None, ())
        assert [list(row) for row in sheet.rows] == [
            ["France", "Germany", "Spain"],
            ["Italy", "Greece", "Portugal"],
        ]
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "| col1 | col2 | col3 |" in markdown
        assert "| France | Germany | Spain |" in markdown

    def test_csv_has_header_true_forces_first_row_as_header(
        self,
        tmp_path: Path,
    ) -> None:
        """``has_header=True`` forces a header even when the heuristic would reject it.

        Empty cells in the first row normally disqualify it; with the
        override, it is still promoted to headers.
        """
        csv_path = tmp_path / "blank_cells.csv"
        # First row has an empty cell — the heuristic would treat all
        # rows as data. The override forces a header anyway.
        _write_csv(csv_path, [["name", "", "city"], ["Alice", "30", "NYC"]])
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        raws = ingestor.discover(tmp_path)
        fragments = ingestor.parse(raws[0], has_header=True)
        sheet = fragments[0].payload
        assert isinstance(sheet, SheetData)
        assert list(sheet.headers or ()) == ["name", "", "city"]
        assert [list(row) for row in sheet.rows] == [["Alice", "30", "NYC"]]

    def test_csv_has_header_none_matches_default_call(
        self,
        tmp_path: Path,
    ) -> None:
        """``has_header=None`` explicit is identical to omitting the kwarg."""
        csv_path = tmp_path / "rows.csv"
        _write_csv(csv_path, [["a", "b"], ["1", "2"]])
        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        raws = ingestor.discover(tmp_path)
        default = ingestor.parse(raws[0])
        explicit_none = ingestor.parse(raws[0], has_header=None)
        default_sheet = default[0].payload
        explicit_sheet = explicit_none[0].payload
        assert isinstance(default_sheet, SheetData)
        assert isinstance(explicit_sheet, SheetData)
        assert default_sheet.headers == explicit_sheet.headers
        assert default_sheet.rows == explicit_sheet.rows

    def test_stub_backend_receives_has_header(self, tmp_path: Path) -> None:
        """``has_header`` propagates from ingestor.parse to backend.read_workbook.

        The stub records every call's ``has_header`` value so we can
        prove the override travels through the ingestor instead of
        being silently dropped on the floor.
        """

        class RecordingBackend:
            """Stub backend that records the has_header it was called with."""

            def __init__(self) -> None:
                """Initialise an empty call-log."""
                self.received: list[bool | None] = []

            def is_available(self) -> bool:
                """Test backend is always available."""
                return True

            def read_workbook(
                self,
                path: Path,
                *,
                has_header: bool | None = None,
            ) -> WorkbookData:
                """Record ``has_header`` and return a canned workbook."""
                self.received.append(has_header)
                return WorkbookData(
                    sheets=(
                        SheetData(
                            name=path.stem,
                            headers=("h1", "h2") if has_header else None,
                            rows=(("v1", "v2"),),
                        ),
                    ),
                )

        path = tmp_path / "data.xlsx"
        _write_xlsx_placeholder(path)
        backend = RecordingBackend()
        ingestor = SpreadsheetIngestor(backend=backend)
        raws = ingestor.discover(tmp_path)
        ingestor.parse(raws[0], has_header=True)
        ingestor.parse(raws[0], has_header=False)
        ingestor.parse(raws[0])
        assert backend.received == [True, False, None]

    def test_xlsx_has_header_false_keeps_first_row_as_data(
        self,
        tmp_path: Path,
    ) -> None:
        """``has_header=False`` flows through the XLSX path as well as CSV.

        Real workbook (via openpyxl) so the override is exercised end
        to end on both spreadsheet formats. Without the override, the
        all-string first row would be promoted to headers by the
        heuristic.
        """
        import openpyxl

        path = tmp_path / "labels.xlsx"
        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        assert worksheet is not None
        worksheet.append(["France", "Germany", "Spain"])
        worksheet.append(["Italy", "Greece", "Portugal"])
        workbook.save(path)

        ingestor = SpreadsheetIngestor(backend=OpenpyxlBackend())
        raws = ingestor.discover(tmp_path)
        fragments = ingestor.parse(raws[0], has_header=False)
        sheet = fragments[0].payload
        assert isinstance(sheet, SheetData)
        assert sheet.headers in (None, ())
        assert [list(row) for row in sheet.rows] == [
            ["France", "Germany", "Spain"],
            ["Italy", "Greece", "Portugal"],
        ]


# ---- OpenpyxlBackend lazy-import ---------------------------------------


class TestOpenpyxlBackend:
    """Default backend respects the optional-dep contract."""

    def test_is_available_returns_false_without_openpyxl(
        self,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Without openpyxl, the backend reports unavailable."""
        import builtins

        monkeypatch.setattr(
            builtins,
            "__import__",
            _make_import_blocker({"openpyxl"}),
        )
        backend = OpenpyxlBackend()
        assert not backend.is_available()

    def test_xlsx_read_raises_without_openpyxl(
        self,
        monkeypatch: pytest.MonkeyPatch,
        tmp_path: Path,
    ) -> None:
        """Reading an .xlsx without openpyxl raises a clear error."""
        import builtins

        monkeypatch.setattr(
            builtins,
            "__import__",
            _make_import_blocker({"openpyxl"}),
        )
        path = tmp_path / "x.xlsx"
        _write_xlsx_placeholder(path)
        backend = OpenpyxlBackend()
        with pytest.raises(OpenpyxlUnavailableError, match="openpyxl"):
            backend.read_workbook(path)


# =======================================================================
# PresentationIngestor
# =======================================================================


class TestPresentationConstants:
    """Public constants exposed by the presentation module."""

    def test_extensions_cover_pptx(self) -> None:
        """``.pptx`` is the only supported extension."""
        assert ".pptx" in PRESENTATION_EXTENSIONS


class TestSlideData:
    """Invariants on the :class:`SlideData` value object."""

    def test_default_notes_is_empty(self) -> None:
        """Slides without speaker notes carry an empty notes string."""
        slide = SlideData(index=1, title="Title", body="Body")
        assert slide.notes == ""


# ---- discover ---------------------------------------------------------


class TestPresentationDiscover:
    """`PresentationIngestor.discover` finds .pptx files."""

    def test_finds_pptx(self, tmp_path: Path) -> None:
        """``.pptx`` files are discovered."""
        _write_pptx_placeholder(tmp_path / "deck.pptx")
        ingestor = PresentationIngestor(backend=StubPresentationBackend())
        raws = ingestor.discover(tmp_path)
        assert [raw.path.name for raw in raws] == ["deck.pptx"]

    def test_ignores_non_pptx_files(self, tmp_path: Path) -> None:
        """Other extensions are skipped."""
        _write_pptx_placeholder(tmp_path / "ok.pptx")
        (tmp_path / "ignored.docx").write_bytes(b"PK")
        ingestor = PresentationIngestor(backend=StubPresentationBackend())
        raws = ingestor.discover(tmp_path)
        assert [raw.path.name for raw in raws] == ["ok.pptx"]


# ---- parse / convert_to_markdown ---------------------------------------


class TestPresentationParse:
    """One :class:`ParsedFragment` per presentation, with all slides inlined."""

    def test_parse_returns_one_fragment(self, tmp_path: Path) -> None:
        """A single presentation file yields a single fragment."""
        path = tmp_path / "deck.pptx"
        _write_pptx_placeholder(path)
        backend = StubPresentationBackend(
            presentations={
                "deck.pptx": PresentationData(
                    title="My Talk",
                    slides=(
                        SlideData(
                            index=1,
                            title="Hello",
                            body="World",
                            notes="speaker note",
                        ),
                        SlideData(
                            index=2,
                            title="Two",
                            body="Bullet 1\nBullet 2",
                        ),
                    ),
                ),
            },
        )
        ingestor = PresentationIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert len(fragments) == 1
        assert fragments[0].metadata["slide_count"] == 2
        assert fragments[0].metadata["title"] == "My Talk"

    def test_parsed_fragment_carries_typed_presentationdata_payload(
        self,
        tmp_path: Path,
    ) -> None:
        """Slide structure rides on a typed :class:`PresentationData` payload.

        Issue #166 — the pre-refactor code dict-ified every slide into
        ``metadata["slides"]`` and ``convert_to_markdown`` re-parsed
        the list of dicts. The typed payload eliminates that round-trip.
        """
        path = tmp_path / "typed.pptx"
        _write_pptx_placeholder(path)
        original = PresentationData(
            title="Typed Deck",
            slides=(SlideData(index=1, title="A", body="B"),),
        )
        backend = StubPresentationBackend(presentations={"typed.pptx": original})
        ingestor = PresentationIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert fragments[0].payload is original
        assert "slides" not in fragments[0].metadata

    def test_markdown_renders_one_section_per_slide(
        self,
        tmp_path: Path,
    ) -> None:
        """The markdown body has ``## Slide 1: Hello`` style sections."""
        path = tmp_path / "deck.pptx"
        _write_pptx_placeholder(path)
        backend = StubPresentationBackend(
            presentations={
                "deck.pptx": PresentationData(
                    title="My Talk",
                    slides=(
                        SlideData(
                            index=1,
                            title="Hello",
                            body="World",
                            notes="speaker note",
                        ),
                        SlideData(index=2, title="Two", body="Bullet"),
                    ),
                ),
            },
        )
        ingestor = PresentationIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "## Slide 1: Hello" in markdown
        assert "## Slide 2: Two" in markdown
        assert "World" in markdown
        assert "Bullet" in markdown
        # Speaker notes get a labelled subsection.
        assert "speaker note" in markdown

    def test_empty_presentation_yields_no_fragment(
        self,
        tmp_path: Path,
    ) -> None:
        """A presentation with no slides is dropped."""
        path = tmp_path / "empty.pptx"
        _write_pptx_placeholder(path)
        backend = StubPresentationBackend(
            presentations={
                "empty.pptx": PresentationData(title=None, slides=()),
            },
        )
        ingestor = PresentationIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert fragments == []


# ---- generate_frontmatter ---------------------------------------------


class TestPresentationFrontmatter:
    """Frontmatter shape for presentation fragments."""

    def test_frontmatter_records_platform_and_slide_count(
        self,
        tmp_path: Path,
    ) -> None:
        """source.platform=presentation plus slide_count and title."""
        from datetime import UTC, datetime

        from creek.ingest.base import ParsedFragment

        fragment = ParsedFragment(
            content="",
            metadata={
                "original_file": str(tmp_path / "deck.pptx"),
                "title": "My Talk",
                "slide_count": 4,
            },
            source_path=str(tmp_path / "deck.pptx"),
            timestamp=datetime(2026, 4, 27, tzinfo=UTC),
        )
        ingestor = PresentationIngestor(backend=StubPresentationBackend())
        fm = ingestor.generate_frontmatter(fragment)
        assert fm["source"]["platform"] == SourcePlatform.PRESENTATION.value
        assert fm["title"] == "My Talk"
        assert fm["slide_count"] == 4


class TestPresentationAuthoredAt:
    """FEAT-031: PPTX core properties ``created`` / ``modified`` → ``authored_at``."""

    def _write_real_pptx(
        self,
        path: Path,
        *,
        created: object | None = None,
        modified: object | None = None,
    ) -> None:
        """Write a real python-pptx-parseable PPTX with the given core properties."""
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
        if created is not None:
            prs.core_properties.created = created  # type: ignore[assignment]
        if modified is not None:
            prs.core_properties.modified = modified  # type: ignore[assignment]
        prs.save(path)

    def test_authored_at_from_pptx_created(self, tmp_path: Path) -> None:
        """The ``created`` core property is promoted to ``authored_at``."""
        from datetime import UTC, datetime

        pptx = tmp_path / "deck.pptx"
        self._write_real_pptx(
            pptx,
            created=datetime(2024, 3, 15, 8, 30, tzinfo=UTC),
        )
        ingestor = PresentationIngestor(backend=PythonPptxBackend())
        fragments = ingestor.parse(ingestor.discover(pptx)[0])
        assert len(fragments) == 1
        authored = fragments[0].metadata["authored_at"]
        assert authored is not None
        assert authored.date() == datetime(2024, 3, 15).date()
        assert authored.tzinfo is not None

    def test_authored_at_in_frontmatter_when_present(self, tmp_path: Path) -> None:
        """Generated frontmatter carries ``authored_at`` as ISO string."""
        from datetime import UTC, datetime

        pptx = tmp_path / "deck.pptx"
        self._write_real_pptx(
            pptx,
            created=datetime(2024, 3, 15, tzinfo=UTC),
        )
        ingestor = PresentationIngestor(backend=PythonPptxBackend())
        fragments = ingestor.parse(ingestor.discover(pptx)[0])
        fm = ingestor.generate_frontmatter(fragments[0])
        assert "authored_at" in fm
        assert fm["authored_at"].startswith("2024-03-15")


# ---- PythonPptxBackend lazy-import -----------------------------------


class TestPythonPptxBackend:
    """Default backend respects the optional-dep contract."""

    def test_is_available_returns_false_without_pptx(
        self,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Without python-pptx, the backend reports unavailable."""
        import builtins

        monkeypatch.setattr(
            builtins,
            "__import__",
            _make_import_blocker({"pptx"}),
        )
        backend = PythonPptxBackend()
        assert not backend.is_available()

    def test_read_raises_without_pptx(
        self,
        monkeypatch: pytest.MonkeyPatch,
        tmp_path: Path,
    ) -> None:
        """Reading a .pptx without python-pptx raises a clear error."""
        import builtins

        monkeypatch.setattr(
            builtins,
            "__import__",
            _make_import_blocker({"pptx"}),
        )
        path = tmp_path / "x.pptx"
        _write_pptx_placeholder(path)
        backend = PythonPptxBackend()
        with pytest.raises(PythonPptxUnavailableError, match="python-pptx"):
            backend.read_presentation(path)

    def test_extract_title_returns_none_when_core_property_is_none(
        self,
    ) -> None:
        """python-pptx sets core_properties.title to None — not the string."""

        class _CoreProps:
            title = None

        class _Prs:
            core_properties = _CoreProps()

        assert PythonPptxBackend._extract_title(_Prs()) is None

    def test_extract_title_strips_whitespace_and_treats_blank_as_none(
        self,
    ) -> None:
        """A title made entirely of whitespace is treated as absent."""

        class _CoreProps:
            title = "   "

        class _Prs:
            core_properties = _CoreProps()

        assert PythonPptxBackend._extract_title(_Prs()) is None

    def test_extract_title_returns_real_title_unchanged(self) -> None:
        """A populated title is returned verbatim (after .strip())."""

        class _CoreProps:
            title = "  My Talk  "

        class _Prs:
            core_properties = _CoreProps()

        assert PythonPptxBackend._extract_title(_Prs()) == "My Talk"


class TestPresentationParseHandlesMissingTitle:
    """Untitled presentations must not surface the literal string ``None``."""

    def test_untitled_pptx_falls_back_to_file_stem(self, tmp_path: Path) -> None:
        """When ``data.title`` is ``None`` the fragment uses the file stem."""
        path = tmp_path / "anonymous.pptx"
        _write_pptx_placeholder(path)
        backend = StubPresentationBackend(
            presentations={
                "anonymous.pptx": PresentationData(
                    title=None,
                    slides=(SlideData(index=1, title="Hi", body="There"),),
                ),
            },
        )
        ingestor = PresentationIngestor(backend=backend)
        fragments = ingestor.parse(ingestor.discover(tmp_path)[0])
        assert fragments[0].metadata["title"] == "anonymous"
        markdown = ingestor.convert_to_markdown(fragments[0])
        assert "# anonymous" in markdown
        # The literal string "None" should never surface as the document title.
        assert "# None" not in markdown


# ---- Module re-exports -----------------------------------------------


class TestModuleExports:
    """Both ingestors are part of the ``creek.ingest`` public surface."""

    def test_spreadsheet_importable_from_package(self) -> None:
        """``from creek.ingest import SpreadsheetIngestor`` works."""
        from creek.ingest import SpreadsheetIngestor as Reexported

        assert Reexported is SpreadsheetIngestor

    def test_presentation_importable_from_package(self) -> None:
        """``from creek.ingest import PresentationIngestor`` works."""
        from creek.ingest import PresentationIngestor as Reexported

        assert Reexported is PresentationIngestor

    def test_protocols_satisfy_structural_typing(self) -> None:
        """Stub backends are structural matches for their Protocols."""
        assert isinstance(StubSpreadsheetBackend(), SpreadsheetBackend)
        assert isinstance(StubPresentationBackend(), PresentationBackend)

    def test_registry_keys_route_correctly(self) -> None:
        """The INGESTOR_REGISTRY keys ``spreadsheet`` and ``presentation`` resolve."""
        from creek.ingest import INGESTOR_REGISTRY

        assert INGESTOR_REGISTRY["spreadsheet"] is SpreadsheetIngestor
        assert INGESTOR_REGISTRY["presentation"] is PresentationIngestor


# =======================================================================
# Per-sheet fragment identity (#1305)
# =======================================================================
#
# ``SpreadsheetIngestor.parse`` emits one ``ParsedFragment`` per non-empty
# sheet, but each one carries ``content=""``, the same hoisted
# ``file_modified_time(raw.path)`` timestamp and the same ``source_path``.
# ``generate_fragment_id(source, timestamp, content)`` therefore mints ONE
# id for all N sheets, and ``VaultWriter._write_model``
# (``VaultWriter._write_model``) takes the lock, calls
# ``_find_existing_locked`` and returns the already-written path for a
# duplicate id — FIRST WRITER WINS. Sheets 2..N are dropped in silence and
# the one surviving file holds the FIRST sheet.
#
# The migration hazard runs the other way: a workbook with exactly ONE
# non-empty sheet (every CSV, every single-sheet XLSX) and every
# presentation must keep the byte-identical id and filename they have
# today, or the fix orphans fragments already sitting in real vaults. The
# ``_PIN_*`` constants below are that guard.


_PINNED_MTIME = 1710500000
"""Fixed epoch mtime for the pin tests, so the hashed timestamp is constant."""

# Measured at HEAD against a RELATIVE source path (``book.xlsx``) whose
# mtime is pinned to ``_PINNED_MTIME``, which makes both hash inputs
# constant across machines and runs. Each value MUST be byte-identical
# after the sheet discriminator lands; a change means an existing
# single-unit fragment has been orphaned.
#
# The filename pins are the TAIL only: ``_compute_base_name`` prefixes
# ``{date}-`` from ``Fragment.created`` (i.e. today), which is not a
# property of this change and would make the pin a calendar flake.
_PIN_XLSX_SINGLE_SHEET_ID = "frag-2ae23100e7a8"
"""Fragment id of a one-sheet ``book.xlsx``; a ``frag-<12hex>`` literal."""

_PIN_XLSX_SINGLE_SHEET_FILENAME_TAIL = "-book.md"
"""Filename tail of a one-sheet ``book.xlsx``; expected ``-book.md``."""

_PIN_CSV_ID = "frag-e2f9d18e46af"
"""Fragment id of ``book.csv``; a ``frag-<12hex>`` literal."""

_PIN_CSV_FILENAME_TAIL = "-book.md"
"""Filename tail of ``book.csv``; expected ``-book.md``."""

_PIN_PRESENTATION_ID = "frag-439246060197"
"""Fragment id of ``deck.pptx``; a ``frag-<12hex>`` literal."""

_PIN_PRESENTATION_FILENAME_TAIL = "-deck.md"
"""Filename tail of ``deck.pptx``; expected ``-deck.md``."""


# ---- Helpers: vault round-trip -----------------------------------------


def _scaffold_vault(tmp_path: Path, name: str = "vault") -> Path:
    """Create the minimum vault tree ``VaultWriter`` refuses to start without.

    ``VaultWriter.__init__`` raises ``FileNotFoundError`` unless both
    ``00-Creek-Meta/`` and ``01-Fragments/`` already exist; the
    per-platform subfolder underneath is created by the writer itself.
    """
    vault = tmp_path / name
    for relpart in ("00-Creek-Meta/Processing-Log", "01-Fragments"):
        (vault / relpart).mkdir(parents=True, exist_ok=True)
    return vault


def _write_real_workbook(path: Path, sheet_names: list[str]) -> None:
    """Write a genuine openpyxl workbook with one sheet per name.

    Every sheet gets a header row plus one data row so none of them is
    dropped by the ``sheet.is_empty`` filter in ``parse``.
    """
    import openpyxl

    workbook = openpyxl.Workbook()
    first = workbook.active
    assert first is not None
    first.title = sheet_names[0]
    for name in sheet_names[1:]:
        workbook.create_sheet(title=name)
    for index, name in enumerate(sheet_names):
        worksheet = workbook[name]
        worksheet.append(["item", "amount"])
        worksheet.append([f"row-{index}", str(index * 10)])
    workbook.save(path)


def _stub_workbook(sheet_names: list[str]) -> WorkbookData:
    """Build canned :class:`WorkbookData` with one non-empty sheet per name.

    Used where openpyxl cannot express the case under test — duplicate
    sheet titles and blank sheet titles are both rejected or silently
    renamed by the real writer, but a hand-built workbook (and a workbook
    produced by some other tool) can carry them.
    """
    return WorkbookData(
        sheets=tuple(
            SheetData(
                name=name,
                headers=("item", "amount"),
                rows=((f"row-{index}", str(index * 10)),),
            )
            for index, name in enumerate(sheet_names)
        ),
    )


def _assemble_all(ingestor: Ingestor, source: Path) -> list[IngestedFragment]:
    """Run the four-stage pipeline and assemble every fragment it yields."""
    result = ingestor.ingest(source)
    assert result.errors == [], result.errors
    return [assemble_ingested_fragment(parsed) for parsed in result.fragments]


def _write_all(vault: Path, assembled: list[IngestedFragment]) -> list[Path]:
    """Write every assembled fragment through the real ``VaultWriter``."""
    writer = VaultWriter(vault_path=vault)
    return [writer.write_fragment(item.fragment, body=item.body) for item in assembled]


def _fragment_files(vault: Path) -> list[Path]:
    """Return every fragment markdown file under ``01-Fragments``, sorted."""
    return sorted((vault / "01-Fragments").rglob("*.md"))


def _title_part(path: Path) -> str:
    """Return a fragment filename's stem with its ``YYYY-MM-DD-`` prefix removed.

    ``_compute_base_name`` builds ``{date}-{sanitized-title}``; the date
    is always the 10 characters of an ISO date, so the remainder is the
    part this change is responsible for.
    """
    return path.stem[11:]


class TestMultiSheetWorkbookIdentity:
    """Every non-empty sheet must survive the write as its own fragment (#1305)."""

    def test_multi_sheet_workbook_writes_one_file_per_sheet(
        self,
        tmp_path: Path,
    ) -> None:
        """A 3-sheet workbook lands as three fragment files, one per sheet.

        Pre-fix symptom (this is what makes the test RED, and it is *not*
        what the issue body claims): all three sheets hash to the same
        ``frag-<12hex>`` because ``parse`` gives each of them
        ``content=""``, the same hoisted ``file_modified_time`` timestamp
        and the same ``source_path``. ``VaultWriter._write_model``
        (``VaultWriter._write_model``) holds the lock, finds the id
        already indexed via ``_find_existing_locked`` and returns the
        existing path without writing — so the FIRST writer wins. One file
        survives and it holds the FIRST sheet (``Budget``); ``Notes`` and
        ``Q3`` are dropped with no error, no warning and no duplicate.
        """
        source = tmp_path / "book.xlsx"
        _write_real_workbook(source, ["Budget", "Notes", "Q3"])
        vault = _scaffold_vault(tmp_path)

        assembled = _assemble_all(SpreadsheetIngestor(), source)
        _write_all(vault, assembled)

        # (a) Three sheets, three distinct identities. No ``frag-`` literal
        # is pinned here: the id hashes the absolute ``tmp_path``, so any
        # literal would differ on every run.
        assert len({item.fragment.id for item in assembled}) == 3

        # (b) Three distinct identities, three files on disk.
        files = _fragment_files(vault)
        assert len(files) == 3

        # (c) Each sheet's rendered heading appears in exactly one file.
        texts = [path.read_text(encoding="utf-8") for path in files]
        for heading in (
            "# book.xlsx — Budget",
            "# book.xlsx — Notes",
            "# book.xlsx — Q3",
        ):
            occurrences = sum(text.count(heading) for text in texts)
            assert occurrences == 1, f"{heading!r} appeared {occurrences} times"

    def test_duplicate_sheet_names_still_get_distinct_ids(
        self,
        tmp_path: Path,
    ) -> None:
        """Sheets named ``Q``, ``Q`` and ``""`` still produce three fragments.

        Excel forbids duplicate sheet titles and openpyxl silently renames
        the second one, so this workbook is built straight from the
        backend's ``WorkbookData`` / ``SheetData`` dataclasses through the
        ``backend`` seam — a file produced by another tool can carry both a
        duplicate title and a blank one. The discriminator must therefore
        be a *normalised, de-duplicated* unit key rather than the raw sheet
        name; a raw-name key collapses the two ``Q`` sheets back into one
        id and the blank-named sheet into an empty discriminator.

        Pre-fix all three sheets share one id and first-writer-wins leaves
        a single file behind.
        """
        source = tmp_path / "book.xlsx"
        _write_xlsx_placeholder(source)
        vault = _scaffold_vault(tmp_path)
        backend = StubSpreadsheetBackend(
            workbooks={"book.xlsx": _stub_workbook(["Q", "Q", ""])},
        )

        assembled = _assemble_all(SpreadsheetIngestor(backend=backend), source)
        assert len(assembled) == 3
        _write_all(vault, assembled)

        assert len({item.fragment.id for item in assembled}) == 3
        assert len(_fragment_files(vault)) == 3

    def test_duplicate_sheet_names_get_distinct_titles_and_headings(
        self,
        tmp_path: Path,
    ) -> None:
        """Two sheets named ``Data`` must be tellable apart *on disk*.

        The sibling test above proves the two sheets get distinct ids and
        two files. That is disambiguation in the index — an operator never
        reads an index. Both user-visible strings, the ``title`` in
        frontmatter and the ``# `` heading in the body, were derived from
        the RAW ``metadata["sheet"]`` name, which is ``Data`` for both
        sheets. So the fix produced two fragments the operator opens and
        cannot distinguish: same title, same heading, same rendered table
        shape, differing only in a ``frag-<12hex>`` id and a ``-1``
        de-collision suffix the writer had to invent because the computed
        filenames collided too.

        The deduplicated ``source_unit`` (``Data`` / ``Data~2``) is the
        string that already resolved this, and is what both surfaces must
        consume. Asserted against file bytes rather than the assembled
        model, because the frontmatter round-trip is where the sibling
        keys ``sheet`` / ``rows`` / ``columns`` are silently dropped
        (#1392) — an in-memory assertion would not have noticed that the
        title is the only per-sheet marker left.
        """
        source = tmp_path / "book.xlsx"
        _write_xlsx_placeholder(source)
        vault = _scaffold_vault(tmp_path)
        backend = StubSpreadsheetBackend(
            workbooks={"book.xlsx": _stub_workbook(["Data", "Data"])},
        )

        assembled = _assemble_all(SpreadsheetIngestor(backend=backend), source)
        _write_all(vault, assembled)

        files = _fragment_files(vault)
        assert len(files) == 2
        loaded = [frontmatter.load(path) for path in files]

        titles = sorted(str(post.metadata["title"]) for post in loaded)
        assert titles == ["book — Data", "book — Data~2"]

        headings = sorted(
            line
            for post in loaded
            for line in post.content.splitlines()
            if line.startswith("# ")
        )
        assert headings == ["# book.xlsx — Data", "# book.xlsx — Data~2"]

        # Distinct titles mean distinct computed filenames, so the writer
        # has no name clash to de-collide. A ``-1`` tail is the on-disk
        # signature of two fragments that both wanted the same name.
        assert sorted(_title_part(path) for path in files) == [
            "book--Data",
            "book--Data2",
        ]

    def test_inserting_a_first_sheet_does_not_re_mint_existing_ids(
        self,
        tmp_path: Path,
    ) -> None:
        """Prepending a sheet leaves the other sheets' ids untouched.

        This is the explicit rejection of the "use an empty discriminator
        for sheet index 0" shortcut. Under that shortcut the workbook's
        first sheet keeps the bare ``source_path`` identity, so inserting a
        new sheet at the front hands the newcomer the id that used to
        belong to the old first sheet. With first-writer-wins in
        ``VaultWriter._write_model``'s duplicate-id early return, the
        newcomer is then SILENTLY DROPPED as a duplicate of a fragment it
        has nothing to do with, and every following sheet shifts one slot
        and overwrites its neighbour's identity.

        Both parses run against the same ``source_path`` string and the
        same pinned mtime, so the timestamp and source inputs to
        ``generate_fragment_id`` are held constant and only the sheet set
        varies. Pre-fix, every sheet in both runs shares one id, so the
        distinctness assertions fail first.
        """
        import os

        source = tmp_path / "book.xlsx"
        _write_xlsx_placeholder(source)
        os.utime(source, (_PINNED_MTIME, _PINNED_MTIME))
        vault = _scaffold_vault(tmp_path)

        before = _assemble_all(
            SpreadsheetIngestor(
                backend=StubSpreadsheetBackend(
                    workbooks={"book.xlsx": _stub_workbook(["Budget", "Notes"])},
                ),
            ),
            source,
        )
        _write_all(vault, before)
        ids_before = {item.fragment.title: item.fragment.id for item in before}
        assert len(set(ids_before.values())) == 2

        after = _assemble_all(
            SpreadsheetIngestor(
                backend=StubSpreadsheetBackend(
                    workbooks={
                        "book.xlsx": _stub_workbook(["Intro", "Budget", "Notes"]),
                    },
                ),
            ),
            source,
        )
        _write_all(vault, after)
        ids_after = {item.fragment.title: item.fragment.id for item in after}
        assert len(set(ids_after.values())) == 3

        # The two pre-existing sheets keep the exact ids they were minted
        # with; only the inserted sheet is new.
        titles_before = sorted(ids_before)
        for title in titles_before:
            assert ids_after[title] == ids_before[title], title
        new_ids = set(ids_after.values()) - set(ids_before.values())
        assert len(new_ids) == 1

        # Persisted state: the insert adds exactly one file and re-writing
        # the two known sheets is a no-op, not a duplicate and not a drop.
        assert len(_fragment_files(vault)) == 3

    def test_re_ingesting_a_workbook_mints_no_new_files(
        self,
        tmp_path: Path,
    ) -> None:
        """A second ingest of an unchanged workbook writes nothing new.

        Idempotency on the default, non-ledgered path: the ids are a pure
        function of ``(source, timestamp, content)``, so the second run
        resolves to the same three fragments and ``_write_model`` returns
        the existing paths. A ``-1`` suffixed stem would mean the writer
        had to invent a filename because two different ids wanted the same
        one.

        Pre-fix this test fails on the first run already — one file, not
        three.
        """
        source = tmp_path / "book.xlsx"
        _write_real_workbook(source, ["Budget", "Notes", "Q3"])
        vault = _scaffold_vault(tmp_path)

        first = _assemble_all(SpreadsheetIngestor(), source)
        _write_all(vault, first)
        assert len(_fragment_files(vault)) == 3

        second = _assemble_all(SpreadsheetIngestor(), source)
        _write_all(vault, second)

        assert {item.fragment.id for item in second} == {
            item.fragment.id for item in first
        }
        files = _fragment_files(vault)
        assert len(files) == 3
        assert [path for path in files if path.stem.endswith("-1")] == []

    def test_multi_sheet_filenames_are_distinct_and_well_formed(
        self,
        tmp_path: Path,
    ) -> None:
        """Each sheet's file is named for its sheet, with no collision suffix.

        ``_sanitize_title`` strips the em-dash and turns the two spaces
        around it into two hyphens, so a title of ``book — Budget``
        becomes the stem tail ``book--Budget`` (the date prefix comes from
        ``Fragment.created`` and is stripped by :func:`_title_part`).

        The blank-sheet-name case is checked in its own vault because
        openpyxl will not create a sheet with an empty title. A blank name
        must normalise to the literal ``sheet`` — an empty discriminator
        would render ``book — `` and sanitise down to a bare ``book``,
        colliding with any other unnamed sheet.

        Pre-fix every sheet is titled ``book`` (``assemble_ingested_fragment``
        falls back to the source stem because the spreadsheet frontmatter
        carries no ``title``), and only one file is written at all.
        """
        source = tmp_path / "book.xlsx"
        _write_real_workbook(source, ["Budget", "Notes", "Q3"])
        vault = _scaffold_vault(tmp_path)
        _write_all(vault, _assemble_all(SpreadsheetIngestor(), source))

        files = _fragment_files(vault)
        assert {_title_part(path) for path in files} == {
            "book--Budget",
            "book--Notes",
            "book--Q3",
        }
        assert [path for path in files if path.stem.endswith("-")] == []
        assert [path for path in files if path.stem.endswith("-1")] == []

        blank_vault = _scaffold_vault(tmp_path, name="vault-blank")
        blank_backend = StubSpreadsheetBackend(
            workbooks={"book.xlsx": _stub_workbook(["Data", ""])},
        )
        blank_source = tmp_path / "blank" / "book.xlsx"
        blank_source.parent.mkdir()
        _write_xlsx_placeholder(blank_source)
        _write_all(
            blank_vault,
            _assemble_all(SpreadsheetIngestor(backend=blank_backend), blank_source),
        )
        blank_files = _fragment_files(blank_vault)
        assert {_title_part(path) for path in blank_files} == {
            "book--Data",
            "book--sheet",
        }
        assert [path for path in blank_files if path.stem.endswith("-")] == []

    def test_sheet_rows_and_columns_survive_to_disk(
        self,
        tmp_path: Path,
    ) -> None:
        """Sheet/rows/columns reach the vault file, and so does the title (#1392).

        **This test is the deliberate inversion of the #1305 pin it
        replaces.** That pin recorded the drop as accepted-for-now and said
        in as many words: "Revisiting the drop is tracked in #1392. If that
        issue decides to surface the fields, this test is the one to update
        — deliberately, not incidentally." #1392 decided to surface them,
        and this is that update. It is not an incidental loosening: the
        assertions were flipped, not deleted, and the two guards that make
        the decision safe live beside it —
        :meth:`test_a_frontmatter_key_outside_the_allowlist_is_still_dropped`
        here, and ``test_a_markdown_fragment_gains_no_dimension_keys`` in
        ``tests/test_ingest_rendered_body_hash.py``, which pins that a
        fragment whose ingestor emitted no dimensions gains no keys.

        The mechanism is the writer's existing ``extra_frontmatter`` seam
        (used today for thread/eddy ``aliases``), gated by an explicit
        allowlist — **not** ``extra="allow"`` on ``Fragment``, which would
        let every ingestor typo become frontmatter, and **not** a nullable
        model field, which ``_write_model``'s ``model_dump(mode="json")``
        would print as ``sheet: null`` on every fragment in the vault.

        Driven through ``run_ingest`` rather than ``_write_all``: the
        passthrough is a property of the real write path, and a helper that
        calls ``write_fragment`` directly steps over the seam under test.

        RED today on the three dimension keys; the ``title`` assertion is
        green and stays green.
        """
        from creek.ingest.pipeline import run_ingest

        source = tmp_path / "book.xlsx"
        _write_real_workbook(source, ["Budget", "Notes", "Q3"])
        vault = _scaffold_vault(tmp_path)
        result = run_ingest(
            ingestor_cls=SpreadsheetIngestor,
            source_type="spreadsheet",
            input_path=source,
            vault_path=vault,
        )
        assert result.errors == [], result.errors

        budget = [
            path
            for path in _fragment_files(vault)
            if "# book.xlsx — Budget" in path.read_text(encoding="utf-8")
        ]
        assert len(budget) == 1
        post = frontmatter.load(str(budget[0]))
        assert post.metadata["sheet"] == "Budget"
        assert post.metadata["rows"] == 1
        assert post.metadata["columns"] == 2
        assert post["title"] == "book — Budget"

    def test_a_frontmatter_key_outside_the_allowlist_is_still_dropped(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """An unlisted key an ingestor emits must never reach the vault (#1392).

        GREEN today and the ratchet on the #1392 decision: the passthrough
        is an explicit allowlist of structured provenance, so a misspelled
        ``sheat`` — or any future key an ingestor grows — is still dropped.
        This is what stops the fix from being implemented as
        ``extra="allow"`` by the back door, where the mistake is invisible
        until it is on disk in a thousand fragments.
        """
        from creek.ingest.pipeline import run_ingest

        real = SpreadsheetIngestor.generate_frontmatter

        def _with_stray_key(
            self: SpreadsheetIngestor,
            fragment: ParsedFragment,
        ) -> dict[str, Any]:
            """Return the real frontmatter plus one key no allowlist names."""
            data = real(self, fragment)
            data["sheat"] = 1
            return data

        monkeypatch.setattr(
            SpreadsheetIngestor, "generate_frontmatter", _with_stray_key
        )

        source = tmp_path / "book.xlsx"
        _write_real_workbook(source, ["Budget"])
        vault = _scaffold_vault(tmp_path)
        result = run_ingest(
            ingestor_cls=SpreadsheetIngestor,
            source_type="spreadsheet",
            input_path=source,
            vault_path=vault,
        )
        assert result.errors == [], result.errors

        files = _fragment_files(vault)
        assert len(files) == 1
        assert "sheat" not in frontmatter.load(str(files[0])).metadata


class TestSingleUnitIdentityIsUnchanged:
    """The migration carve-out: one-unit sources keep their exact identity.

    Every test in this class is GREEN BEFORE AND AFTER the fix, by
    construction — that is the point. A workbook with exactly one non-empty
    sheet (which is every CSV and every single-sheet XLSX) and every
    presentation must keep the byte-identical ``frag-<12hex>`` and the
    byte-identical filename they have today. If any of these goes red, the
    change has orphaned fragments that already exist in real vaults, and no
    amount of correctness on the multi-sheet path pays for that.

    Each test pins against a RELATIVE source path under a chdir'd
    ``tmp_path`` and an mtime forced to :data:`_PINNED_MTIME`, so both
    inputs ``generate_fragment_id`` hashes are constants rather than
    per-run values.
    """

    def test_single_sheet_xlsx_id_is_unchanged_by_the_sheet_discriminator(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """A one-sheet XLSX keeps its measured id and filename.

        Green before and after. The carve-out under test: the sheet
        discriminator is set only when a workbook yields two or more
        non-empty sheets, so a single-sheet workbook keeps
        ``source_unit=None``, hashes the bare ``source_path``, and titles
        itself from the plain file stem.
        """
        import os
        from pathlib import Path

        monkeypatch.chdir(tmp_path)
        source = Path("book.xlsx")
        _write_real_workbook(source, ["Sheet1"])
        os.utime(source, (_PINNED_MTIME, _PINNED_MTIME))
        vault = _scaffold_vault(tmp_path)

        assembled = _assemble_all(SpreadsheetIngestor(), source)
        assert len(assembled) == 1
        assert assembled[0].fragment.id == _PIN_XLSX_SINGLE_SHEET_ID

        written = _write_all(vault, assembled)
        assert len(_fragment_files(vault)) == 1
        assert written[0].name.endswith(_PIN_XLSX_SINGLE_SHEET_FILENAME_TAIL)

    def test_single_sheet_csv_id_is_unchanged_by_the_sheet_discriminator(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """A CSV keeps its measured id and filename.

        Green before and after, and load-bearing: ``_csv_text_to_workbook``
        (creek/ingest/spreadsheets.py ~336) names a CSV's lone sheet after
        the file stem, so a discriminator applied unconditionally — even a
        name-based one — would move the id of every CSV fragment in every
        vault. This test is what forbids that.
        """
        import os
        from pathlib import Path

        monkeypatch.chdir(tmp_path)
        source = Path("book.csv")
        source.write_text("a,b\n1,2\n3,4\n", encoding="utf-8")
        os.utime(source, (_PINNED_MTIME, _PINNED_MTIME))
        vault = _scaffold_vault(tmp_path)

        assembled = _assemble_all(SpreadsheetIngestor(), source)
        assert len(assembled) == 1
        assert assembled[0].fragment.id == _PIN_CSV_ID

        written = _write_all(vault, assembled)
        assert len(_fragment_files(vault)) == 1
        assert written[0].name.endswith(_PIN_CSV_FILENAME_TAIL)

    def test_presentation_fragment_id_is_unchanged_by_the_sheet_discriminator(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """A presentation keeps its measured id and filename.

        Green before and after. ``PresentationIngestor.parse``
        (creek/ingest/presentations.py ~310) has the identical
        ``content=""`` shape as the spreadsheet ingestor, so it is the
        obvious place for a sheet-discriminator refactor to leak into. It
        must be provably untouched: a deck is one unit and stays one unit.
        """
        import os
        from pathlib import Path

        monkeypatch.chdir(tmp_path)
        source = Path("deck.pptx")
        _write_pptx_placeholder(source)
        os.utime(source, (_PINNED_MTIME, _PINNED_MTIME))
        vault = _scaffold_vault(tmp_path)
        backend = StubPresentationBackend(
            presentations={
                "deck.pptx": PresentationData(
                    title=None,
                    slides=(SlideData(index=1, title="Hello", body="World"),),
                ),
            },
        )

        assembled = _assemble_all(PresentationIngestor(backend=backend), source)
        assert len(assembled) == 1
        assert assembled[0].fragment.id == _PIN_PRESENTATION_ID

        written = _write_all(vault, assembled)
        assert len(_fragment_files(vault)) == 1
        assert written[0].name.endswith(_PIN_PRESENTATION_FILENAME_TAIL)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
