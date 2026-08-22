"""The ledger must hash the body the vault actually holds (#1393/#1392/#1482).

Three defects share one root and are pinned together here.

**#1393 — silent data loss.** ``SpreadsheetIngestor.parse`` and
``PresentationIngestor.parse`` both return ``ParsedFragment(content="")`` and
defer the real body to ``convert_to_markdown``. ``record_in_ledger`` hashes
``parsed.content``, so every workbook and every deck records
``sha256(b"")`` — ``e3b0c442…`` — forever. An edited re-upload therefore
compares equal to its own predecessor, takes the ``unchanged`` branch, and the
operator's edit is never written. The run reports success.

**#1392 — provenance dropped at the door.** The spreadsheet ingestor emits
``sheet``/``rows``/``columns`` on its frontmatter dict; ``Fragment`` leaves
pydantic's default ``extra="ignore"``, so all three are discarded before
anything reaches disk. The #1305 pin recorded that drop as accepted-for-now
and named this issue as the place to decide it. Both entry points are pinned:
``creek ingest`` writes them through the assembler, and ``creek process`` has
its own drop — the classification stage re-wraps every ``IngestedFragment``
and used to leave the passthrough dict behind — so the two are asserted
against each other over one source directory.

**#1482 — the summary that adds up to nothing.** ``creek ingest`` prints
created/updated/tombed/skipped and omits ``unchanged``, so a run that writes N
fragments can print all zeros above ``Ingested N fragment(s).``

Everything here drives the REAL surfaces — ``creek.upload`` (the only path on
which spreadsheets and presentations are ledgered at all) and ``creek.cli``'s
``_run_ingest`` — because the defect lives in the seam between the ingestor,
the ledger and the writer, and a test that assembles fragments by hand steps
over exactly that seam.

Every "an edit is now seen" test carries its twin: an UNMODIFIED re-send still
reports ``unchanged``. Without the twin, the first test is satisfiable by
breaking idempotency — trading a silent no-op for a spurious rewrite of every
workbook on every run.
"""

from __future__ import annotations

import base64
import re
from typing import TYPE_CHECKING, Any

import frontmatter
import pytest

from creek.cli import _run_ingest
from creek.config import CreekConfig
from creek.ingest import INGESTOR_REGISTRY, route_to_ingestor
from creek.ingest.ledger import SourceLedger
from creek.ingest.pipeline import run_ingest
from creek.pipeline import Pipeline
from creek_mcp.tier_ceiling import TierCeiling
from creek_mcp.tools.upload import UPLOAD_LEDGER_SOURCE, upload_tool

if TYPE_CHECKING:
    from pathlib import Path

_TS = "2026-06-20T10:00:00+00:00"
"""Fixed upload timestamp, so no assertion here depends on the wall clock."""

_EMPTY_HASH = SourceLedger.content_hash("")
"""Derived, never pasted: the digest a body-less ledger row records today."""

_ORIGINAL_CELL = "cell-value-original-1393"
_EDITED_CELL = "cell-value-edited-1393"
_ORIGINAL_SLIDE = "slide-body-original-1393"
_EDITED_SLIDE = "slide-body-edited-1393"


def _vault(tmp_path: Path, name: str = "vault") -> Path:
    """Create the minimum vault layout the writer, ledger and audit log need."""
    vault = tmp_path / name
    for sub in (
        "00-Creek-Meta/State",
        "00-Creek-Meta/Processing-Log",
        "00-Creek-Meta/audit",
        "01-Fragments",
    ):
        (vault / sub).mkdir(parents=True, exist_ok=True)
    return vault


def _fragments(vault: Path) -> list[Path]:
    """Return every fragment markdown file under ``01-Fragments``, sorted."""
    return sorted((vault / "01-Fragments").rglob("*.md"))


def _fragment_text(vault: Path) -> str:
    """Return every fragment file's text concatenated, for a presence check."""
    return "\n".join(path.read_text(encoding="utf-8") for path in _fragments(vault))


def _b64(payload: bytes) -> str:
    """Return base64 of *payload*, derived at runtime, never hardcoded."""
    return base64.b64encode(payload).decode("ascii")


def _workbook_bytes(tmp_path: Path, value: str) -> bytes:
    """Build a genuine single-sheet ``.xlsx`` whose one data cell is *value*.

    Single-sheet on purpose: the multi-sheet path already has its own suite
    (#1305), and a one-sheet workbook keeps a bare ``origin_key``, so nothing
    here can pass by accident of the sub-unit scheme.

    Args:
        tmp_path: Directory to build the workbook in.
        value: The cell value that distinguishes one revision from the next.

    Returns:
        The workbook's bytes.
    """
    from openpyxl import Workbook

    path = tmp_path / f"book-{value}.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Budget"
    sheet["A1"] = "label"
    sheet["B1"] = "value"
    sheet["A2"] = "rent"
    sheet["B2"] = value
    workbook.save(path)
    return path.read_bytes()


def _deck_bytes(tmp_path: Path, body: str) -> bytes:
    """Build a genuine one-slide ``.pptx`` whose slide body is *body*.

    Args:
        tmp_path: Directory to build the deck in.
        body: The slide's body text, distinguishing one revision from the next.

    Returns:
        The presentation's bytes.
    """
    from pptx import Presentation

    path = tmp_path / f"deck-{body}.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly"
    slide.placeholders[1].text = body
    presentation.save(path)
    return path.read_bytes()


def _send(vault: Path, *, filename: str, payload: bytes, external_id: str) -> Any:
    """Upload *payload* under a stable ``external_id`` at tier ``open``.

    Args:
        vault: Vault root.
        filename: The caller's filename; only its extension is trusted.
        payload: The document bytes.
        external_id: The caller's idempotency key, held constant across a
            re-send so the ledger resolves both runs to one source unit.

    Returns:
        The tool's response mapping.
    """
    return upload_tool(
        vault_path=vault,
        filename=filename,
        content_base64=_b64(payload),
        external_id=external_id,
        tier="open",
        timestamp=_TS,
        privacy_tier_ceiling=TierCeiling.OPEN,
    )


def _origin_key(vault: Path) -> str:
    """Return the single fragment's ``source.origin_key``.

    Args:
        vault: Vault root, expected to hold exactly one fragment.

    Returns:
        The fragment's origin key — the ledger's own key for this unit.
    """
    fragments = _fragments(vault)
    assert len(fragments) == 1, [path.name for path in fragments]
    source = frontmatter.load(str(fragments[0])).metadata["source"]
    assert isinstance(source, dict)
    return str(source["origin_key"])


def _upload_hash(vault: Path) -> str:
    """Return the upload ledger's recorded ``content_hash`` for the one unit.

    Args:
        vault: Vault root, expected to hold exactly one fragment.

    Returns:
        The recorded SHA-256 hex digest.
    """
    ledger = SourceLedger.load(vault, source=UPLOAD_LEDGER_SOURCE)
    record = ledger.get(_origin_key(vault))
    assert record is not None
    return record.content_hash


# ---- #1393: an edited document must be seen ------------------------------


def test_pptx_routes_to_the_presentation_ingestor() -> None:
    """``.pptx`` is reachable through ``creek.upload``; #1393's escape hatch is void.

    Issue #1393 offers "presentations may not be ledger-reachable — confirm
    before fixing". Since #1526 they are: ``_EXTENSION_ROUTES`` maps ``.pptx``
    to ``presentation`` and the upload tool always passes
    ``ledger_source=UPLOAD_LEDGER_SOURCE``. Pinned so the hatch cannot be
    reinstated by a later reader of the issue.

    GREEN today — a premise pin, not a red test.
    """
    from pathlib import Path as _Path

    assert route_to_ingestor(_Path("deck.pptx")) == "presentation"
    assert route_to_ingestor(_Path("book.xlsx")) == "spreadsheet"


def test_edited_workbook_resend_is_updated_and_lands_on_disk(tmp_path: Path) -> None:
    """A re-uploaded workbook with an edited cell must update the fragment.

    RED today, and the failure is the whole issue: the ledger holds
    ``sha256("")`` for run one, run two hashes ``parsed.content`` — still
    ``""`` — the two compare equal, and ``write_fragment_idempotent`` returns
    ``unchanged``. The edit is discarded, permanently: nothing ever revisits
    an unchanged unit.
    """
    vault = _vault(tmp_path)

    first = _send(
        vault,
        filename="budget.xlsx",
        payload=_workbook_bytes(tmp_path, _ORIGINAL_CELL),
        external_id="u-edit-xlsx",
    )
    assert first["status"] == "ok"
    assert first["action"] == "created"
    assert _ORIGINAL_CELL in _fragment_text(vault)

    second = _send(
        vault,
        filename="budget.xlsx",
        payload=_workbook_bytes(tmp_path, _EDITED_CELL),
        external_id="u-edit-xlsx",
    )

    assert second["status"] == "ok"
    assert second["action"] == "updated"
    assert len(_fragments(vault)) == 1
    body = _fragment_text(vault)
    assert _EDITED_CELL in body
    assert _ORIGINAL_CELL not in body


def test_unmodified_workbook_resend_still_reports_unchanged(tmp_path: Path) -> None:
    """The twin: identical bytes must remain a true no-op.

    GREEN today and required to stay green. Without it, the red test above is
    satisfiable by making every unit permanently "changed" — which would
    rewrite every workbook in the vault on every run, re-deriving privacy
    tiers and flagging re-classification for content nobody touched. That is a
    worse defect than the one being fixed, and it would ship looking green.
    """
    vault = _vault(tmp_path)
    payload = _workbook_bytes(tmp_path, _ORIGINAL_CELL)

    first = _send(
        vault, filename="budget.xlsx", payload=payload, external_id="u-same-xlsx"
    )
    key_after_first = _origin_key(vault)
    second = _send(
        vault, filename="budget.xlsx", payload=payload, external_id="u-same-xlsx"
    )

    assert first["action"] == "created"
    assert second["action"] == "unchanged"
    assert len(_fragments(vault)) == 1
    # The origin key is the RTBF purge key; a key that moves on re-ingest
    # drops the old fragment out of every purge the vault offers.
    assert _origin_key(vault) == key_after_first


def test_edited_presentation_resend_is_updated_and_lands_on_disk(
    tmp_path: Path,
) -> None:
    """The same defect, on the deck path #1393 suspected was unreachable.

    RED today: measured ``unchanged`` with the original slide body still on
    disk. ``PresentationIngestor.parse`` carries the identical
    ``content=""`` shape, so it records the identical empty-string digest.
    """
    vault = _vault(tmp_path)

    first = _send(
        vault,
        filename="deck.pptx",
        payload=_deck_bytes(tmp_path, _ORIGINAL_SLIDE),
        external_id="u-edit-pptx",
    )
    assert first["status"] == "ok"
    assert first["action"] == "created"
    assert _ORIGINAL_SLIDE in _fragment_text(vault)

    second = _send(
        vault,
        filename="deck.pptx",
        payload=_deck_bytes(tmp_path, _EDITED_SLIDE),
        external_id="u-edit-pptx",
    )

    assert second["status"] == "ok"
    assert second["action"] == "updated"
    assert len(_fragments(vault)) == 1
    body = _fragment_text(vault)
    assert _EDITED_SLIDE in body
    assert _ORIGINAL_SLIDE not in body


def test_unmodified_presentation_resend_still_reports_unchanged(
    tmp_path: Path,
) -> None:
    """The deck's twin: identical bytes must remain a true no-op.

    GREEN today, for the same reason and with the same job as the workbook
    twin above.
    """
    vault = _vault(tmp_path)
    payload = _deck_bytes(tmp_path, _ORIGINAL_SLIDE)

    first = _send(vault, filename="deck.pptx", payload=payload, external_id="u-same-p")
    key_after_first = _origin_key(vault)
    second = _send(vault, filename="deck.pptx", payload=payload, external_id="u-same-p")

    assert first["action"] == "created"
    assert second["action"] == "unchanged"
    assert len(_fragments(vault)) == 1
    assert _origin_key(vault) == key_after_first


@pytest.mark.parametrize(
    ("filename", "builder", "original", "edited", "external_id"),
    [
        ("budget.xlsx", "workbook", _ORIGINAL_CELL, _EDITED_CELL, "u-hash-xlsx"),
        ("deck.pptx", "deck", _ORIGINAL_SLIDE, _EDITED_SLIDE, "u-hash-pptx"),
    ],
)
def test_ledger_hash_is_a_body_hash_not_the_empty_digest(
    tmp_path: Path,
    filename: str,
    builder: str,
    original: str,
    edited: str,
    external_id: str,
) -> None:
    """The recorded hash must describe the body, and must move when it moves.

    RED today on both arms: the row literally reads
    ``e3b0c44298fc1c149afbf4c8996fb92427ae41e4649b934ca495991b7852b855``
    before and after the edit. This is the direct anti-vacuity assertion for
    the fix — the ``updated`` tests above could in principle be satisfied by
    some other comparison; this one names the stored value.

    Deliberately whitespace-agnostic: it pins that the digest is *derived from
    the body*, not which trailing-newline convention the hash uses. The exact
    convention is pinned separately, against markdown, by
    :func:`test_markdown_ledger_hash_still_matches_the_stored_body`.

    Args:
        tmp_path: Pytest temp dir.
        filename: Caller-supplied name whose extension picks the ingestor.
        builder: ``"workbook"`` or ``"deck"``, naming the fixture builder.
        original: Marker text for the first revision.
        edited: Marker text for the second revision.
        external_id: Idempotency key, held constant across both sends.
    """
    build = _workbook_bytes if builder == "workbook" else _deck_bytes
    vault = _vault(tmp_path)

    _send(
        vault,
        filename=filename,
        payload=build(tmp_path, original),
        external_id=external_id,
    )
    first_hash = _upload_hash(vault)
    _send(
        vault,
        filename=filename,
        payload=build(tmp_path, edited),
        external_id=external_id,
    )
    second_hash = _upload_hash(vault)

    assert first_hash != _EMPTY_HASH
    assert second_hash != _EMPTY_HASH
    assert second_hash != first_hash


# ---- HAZARD 2: the third hash site, reachable only via --incremental -----
#
# `should_skip_unit` hashes `parsed.content` too, and it is a SEPARATE call
# site from the two the upload path exercises. A sweep that fixes the record
# and the compare but leaves the incremental filter hashing "" makes every
# workbook permanently "changed": the filter's hash never matches the row's,
# so `creek ingest --incremental` re-writes the whole corpus on every run.
# These two tests are the only thing standing between that and a green gate.


def _ingest_workbook(vault: Path, source: Path, *, incremental: bool = False) -> Any:
    """Run the ledgered spreadsheet ingest over *source*.

    ``ledger_source`` is passed explicitly because ``spreadsheet`` is not in
    ``LEDGERED_SOURCES``; without it there is no ledger at all and every
    assertion about hashes and skips would be vacuously true.

    Args:
        vault: Vault root.
        source: The workbook (or directory) to ingest.
        incremental: Whether to run the ledger-driven incremental filter.

    Returns:
        The :class:`~creek.ingest.pipeline.IngestRunResult`.
    """
    return run_ingest(
        ingestor_cls=INGESTOR_REGISTRY["spreadsheet"],
        source_type="spreadsheet",
        input_path=source,
        vault_path=vault,
        ledger_source=UPLOAD_LEDGER_SOURCE,
        incremental=incremental,
    )


def test_incremental_run_skips_an_unmodified_workbook(tmp_path: Path) -> None:
    """``--incremental`` must still skip a workbook nobody edited.

    GREEN today — and it is the test that catches a two-of-three sweep. Fix
    the record site and the compare site while leaving ``should_skip_unit``
    hashing ``parsed.content`` and this goes red with ``skipped == 0``,
    because the filter's ``""`` digest can never again match the row's body
    digest. That is the state HAZARD 2 warns is worse than not fixing at all.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    (source / "book.xlsx").write_bytes(_workbook_bytes(tmp_path, _ORIGINAL_CELL))

    first = _ingest_workbook(vault, source)
    second = _ingest_workbook(vault, source, incremental=True)

    assert first.created == 1
    assert second.skipped == 1
    assert second.written == 0


def test_incremental_run_does_not_skip_an_edited_workbook(tmp_path: Path) -> None:
    """``--incremental`` must not skip a workbook whose cell changed.

    RED today: the filter hashes ``""`` on both runs, matches the row, and
    skips — so the incremental surface loses the edit exactly as the upload
    surface does.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    workbook = source / "book.xlsx"
    workbook.write_bytes(_workbook_bytes(tmp_path, _ORIGINAL_CELL))

    _ingest_workbook(vault, source)
    workbook.write_bytes(_workbook_bytes(tmp_path, _EDITED_CELL))
    second = _ingest_workbook(vault, source, incremental=True)

    assert second.skipped == 0
    assert second.updated == 1
    assert _EDITED_CELL in _fragment_text(vault)


# ---- The zero-churn pin: markdown must not move at all ------------------


def _write_journal(source: Path, body: str) -> Path:
    """Write a dated journal entry under *source* and return its path.

    Args:
        source: Directory to write into.
        body: The entry's body text.

    Returns:
        The entry's path.
    """
    entry = source / "2026-06-26.md"
    entry.write_text(f"---\ndate: 2026-06-26\n---\n{body}\n", encoding="utf-8")
    return entry


def _ingest_markdown(
    vault: Path, source: Path, *, print_summary: bool = False
) -> tuple[int, list[str], int]:
    """Run the markdown ingestor through the CLI helper.

    Args:
        vault: Vault root.
        source: Source file or directory.
        print_summary: Whether to print the operator summary line.

    Returns:
        ``(written, errors, discovered)``.
    """
    return _run_ingest(
        ingestor_cls=INGESTOR_REGISTRY["markdown"],
        source_type="markdown",
        input_path=source,
        vault_path=vault,
        print_summary=print_summary,
    )


def test_markdown_ledger_hash_still_matches_the_stored_body(tmp_path: Path) -> None:
    """Markdown's recorded hash is the hash of the body on disk. Unchanged.

    GREEN before and after, by construction — that is the point.
    ``MarkdownIngestor.convert_to_markdown`` is the identity function, so
    hashing the rendered body is hashing what was already hashed. Markdown is
    the default-ledgered type, the ``creek sync`` journal population and the
    entire ``creek ingest --pin-source-ids`` population; if this goes red the
    change has desynchronised the #1329 migration
    (``pin_ids._stored_body_hash``, which already hashes ``post.content``)
    from the runner, and re-ingesting a pinned vault would report ``updated``
    for every fragment in it.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    _write_journal(source, "A quiet entry about slow mornings.")

    written, errors, _ = _ingest_markdown(vault, source)
    assert (written, errors) == (1, [])

    post = frontmatter.load(str(_fragments(vault)[0]))
    ledger = SourceLedger.load(vault, source="markdown")
    record = ledger.get(str(post.metadata["source"]["origin_key"]))
    assert record is not None
    assert record.content_hash == SourceLedger.content_hash(post.content)
    assert record.content_hash != _EMPTY_HASH


# ---- #1392: sheet/rows/columns are provenance an automated consumer needs -


def test_workbook_frontmatter_carries_sheet_rows_and_columns(tmp_path: Path) -> None:
    """The three dimension keys must reach the vault file (#1392).

    RED today. ``SpreadsheetIngestor.generate_frontmatter`` emits ``sheet``,
    ``rows`` and ``columns``; ``Fragment`` leaves pydantic's default
    ``extra="ignore"``, so ``Fragment.model_validate`` discards them and
    ``_write_model`` serialises model fields only. #1305 pinned that drop as
    accepted-for-now and named this issue as where to decide it; this is the
    decision, taken through the writer's existing ``extra_frontmatter`` seam
    so no nullable model field is added and ``sheet: null`` stays structurally
    unreachable.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    (source / "book.xlsx").write_bytes(_workbook_bytes(tmp_path, _ORIGINAL_CELL))

    _ingest_workbook(vault, source)

    post = frontmatter.load(str(_fragments(vault)[0]))
    assert post.metadata["sheet"] == "Budget"
    assert post.metadata["rows"] == 1
    assert post.metadata["columns"] == 2


def test_dimension_keys_survive_the_edited_rewrite(tmp_path: Path) -> None:
    """The dimension keys outlive the ``updated`` branch that #1393 unlocked.

    The one place the two fixes meet, and the only unguarded seam between
    them. #1393 is what makes an edited workbook reach ``update_fragment`` at
    all — before it, the run took the ``unchanged`` branch and this path was
    dead. And ``update_fragment`` is the single write branch that is NOT
    handed ``extra_frontmatter``: it reloads the post and rewrites only
    ``post.content``, so ``sheet``/``rows``/``columns`` survive because they
    are already on disk, not because they were passed in again.

    That makes the durability an inherited property rather than a stated one,
    which is precisely the kind of thing that breaks silently. Were
    ``update_fragment`` ever changed to re-serialise from the model — a
    reasonable-looking refactor — every edited workbook would quietly shed its
    provenance on the *second* ingest, with the first still looking correct.
    Asserted together with the edit landing, so neither half can regress
    without the other being noticed.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    workbook = source / "book.xlsx"
    workbook.write_bytes(_workbook_bytes(tmp_path, _ORIGINAL_CELL))

    _ingest_workbook(vault, source)
    workbook.write_bytes(_workbook_bytes(tmp_path, _EDITED_CELL))
    second = _ingest_workbook(vault, source)

    assert second.updated == 1
    post = frontmatter.load(str(_fragments(vault)[0]))
    assert _EDITED_CELL in post.content
    assert post.metadata["sheet"] == "Budget"
    assert post.metadata["rows"] == 1
    assert post.metadata["columns"] == 2


def test_a_markdown_fragment_gains_no_dimension_keys(tmp_path: Path) -> None:
    """A fragment whose ingestor emitted no dimensions gains no keys.

    GREEN today and the anti-``sheet: null`` guard: the passthrough must be
    keyed on PRESENCE in the ingestor's own frontmatter dict, never on a
    nullable field added to ``Fragment`` or ``FragmentSource``. ``_write_model``
    dumps with ``model_dump(mode="json")`` and no ``exclude_none``, so a
    nullable field would print ``sheet: null`` on every fragment from every
    ingestor — the trap this pins shut. The exact key sets are asserted, not
    just the three names, so any future nullable field fails here loudly.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    _write_journal(source, "An entry with no sheets and no slides.")

    _ingest_markdown(vault, source)

    post = frontmatter.load(str(_fragments(vault)[0]))
    assert {"sheet", "rows", "columns"}.isdisjoint(post.metadata)
    source_block = post.metadata["source"]
    assert isinstance(source_block, dict)
    assert {"sheet", "rows", "columns"}.isdisjoint(source_block)


def _process(vault: Path, source: Path) -> Any:
    """Run the ``creek process`` pipeline over *source* into *vault*.

    ``no_llm=True`` keeps it hermetic: without it the run reaches the LLM
    classifier on any machine with ollama up or a provider key exported,
    which is network egress from a test about frontmatter keys.

    Args:
        vault: Vault root to write into.
        source: Directory of source files to process.

    Returns:
        The :class:`~creek.pipeline.PipelineResult`.
    """
    return Pipeline(config=CreekConfig(), no_llm=True).run(
        source_path=source, vault_path=vault
    )


def test_creek_process_writes_the_dimension_keys_creek_ingest_writes(
    tmp_path: Path,
) -> None:
    """The two entry points must frontmatter the same workbook the same way.

    RED before the ``extra_frontmatter`` forwarding in
    ``Pipeline._run_classification``: ``creek ingest`` wrote
    ``sheet``/``rows``/``columns``, and ``creek process`` — over the byte-identical
    workbook — wrote none of them. The classification stage re-wrapped each
    ``IngestedFragment`` around its classified ``Fragment`` and dropped the
    passthrough dict on the way, so the write stage below it forwarded ``{}``
    every time and its forwarding was dead code.

    Both surfaces read the SAME source directory, so nothing here can pass by
    the two runs having seen different bytes; and the ingest side is asserted
    to the literal values as well, so this cannot go green by both surfaces
    dropping the keys together.
    """
    source = tmp_path / "src"
    source.mkdir()
    (source / "book.xlsx").write_bytes(_workbook_bytes(tmp_path, _ORIGINAL_CELL))
    ingest_vault = _vault(tmp_path, "vault-ingest")
    process_vault = _vault(tmp_path, "vault-process")

    _ingest_workbook(ingest_vault, source)
    result = _process(process_vault, source)

    assert result.errors == []
    ingested = frontmatter.load(str(_fragments(ingest_vault)[0])).metadata
    processed = frontmatter.load(str(_fragments(process_vault)[0])).metadata
    assert ingested["sheet"] == "Budget"
    assert ingested["rows"] == 1
    assert ingested["columns"] == 2
    for key in ("sheet", "rows", "columns"):
        assert processed[key] == ingested[key], key


def test_creek_process_gives_a_markdown_fragment_no_dimension_keys(
    tmp_path: Path,
) -> None:
    """The presence-keying survives on the process path too.

    The twin of the parity test above, and the reason that one cannot be
    satisfied by stamping the three keys unconditionally: a journal entry
    processed by the same pipeline must gain no dimension keys at all,
    because its ingestor emitted none. The exact key sets are asserted so a
    future nullable ``Fragment`` field — the ``sheet: null``-on-everything
    trap ``PASSTHROUGH_FRONTMATTER_KEYS`` exists to avoid — fails here loudly.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    _write_journal(source, "A processed entry with no sheets and no slides.")

    result = _process(vault, source)

    assert result.errors == []
    post = frontmatter.load(str(_fragments(vault)[0]))
    assert {"sheet", "rows", "columns"}.isdisjoint(post.metadata)
    source_block = post.metadata["source"]
    assert isinstance(source_block, dict)
    assert {"sheet", "rows", "columns"}.isdisjoint(source_block)


# ---- #1482: a summary that accounts for every fragment it wrote ----------


def _summary_counts(captured: str) -> dict[str, int]:
    """Parse the ``Ingest summary:`` line into ``{label: count}``.

    Whitespace is normalised first so a rich console wrap cannot split a
    ``"N label"`` pair and silently drop it from the parse.

    Args:
        captured: Everything the run printed to stdout.

    Returns:
        The summary line's counts, keyed by their printed labels.
    """
    flattened = " ".join(captured.split())
    assert "Ingest summary:" in flattened, flattened
    line = flattened.split("Ingest summary:", 1)[1]
    return {label: int(count) for count, label in re.findall(r"(\d+) ([a-z]+)", line)}


def test_ingest_summary_names_the_unchanged_count(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    """A no-op re-run must say so, instead of printing four zeros.

    RED today: the summary is created/updated/tombed/skipped, so the one
    number that explains the run is the one it does not print.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    _write_journal(source, "An entry that will be re-ingested untouched.")

    _ingest_markdown(vault, source, print_summary=True)
    assert _summary_counts(capsys.readouterr().out)["created"] == 1

    _ingest_markdown(vault, source, print_summary=True)
    counts = _summary_counts(capsys.readouterr().out)

    assert counts["unchanged"] == 1
    assert counts["created"] == 0
    assert counts["updated"] == 0


def test_summary_accounts_for_every_fragment_the_run_wrote(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    """The printed numbers must sum to the ``Ingested N fragment(s).`` number.

    RED today, and this is #1482's verbatim reproduction. The trigger is not
    an empty ledger — that takes the ``record is None`` branch and reports
    ``created``. It is a ledger row whose hash still matches while the
    fragment FILE is gone, which is exactly what ``creek purge vault`` leaves
    behind: the unchanged branch calls ``write_fragment``, which recreates the
    file, and returns ``unchanged``. The run writes a fragment, the summary
    prints all zeros, and the next line says ``Ingested 1 fragment(s).``

    ``written`` is what ``creek ingest`` prints in that line (``cli.py``
    ~2106), so asserting against the return value asserts against the printed
    claim.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    _write_journal(source, "An entry that survives its own vault purge.")

    _ingest_markdown(vault, source, print_summary=True)
    capsys.readouterr()
    for path in _fragments(vault):
        path.unlink()

    written, _, _ = _ingest_markdown(vault, source, print_summary=True)
    counts = _summary_counts(capsys.readouterr().out)

    assert written == 1
    assert len(_fragments(vault)) == 1
    accounted = counts["created"] + counts["updated"] + counts.get("unchanged", 0)
    assert accounted == written


def test_run_result_written_equals_created_plus_updated_plus_unchanged(
    tmp_path: Path,
) -> None:
    """The invariant the summary is supposed to expose, over a mixed run.

    GREEN today — the counting has always been right; only the display drops
    ``unchanged``. Pinned as the contract the summary must not be allowed to
    contradict, and as the ratchet that stops a later "simplification" of the
    tally from making the honest summary dishonest again.
    """
    vault = _vault(tmp_path)
    source = tmp_path / "src"
    source.mkdir()
    stable = source / "stable.md"
    stable.write_text("---\ndate: 2026-06-01\n---\nUnchanged across runs.\n", "utf-8")
    edited = source / "edited.md"
    edited.write_text("---\ndate: 2026-06-02\n---\nThe first version.\n", "utf-8")

    _ingest_markdown(vault, source)
    edited.write_text(
        "---\ndate: 2026-06-02\n---\nA wholly different second version.\n", "utf-8"
    )
    fresh = source / "fresh.md"
    fresh.write_text("---\ndate: 2026-06-03\n---\nBrand new this run.\n", "utf-8")

    result = run_ingest(
        ingestor_cls=INGESTOR_REGISTRY["markdown"],
        source_type="markdown",
        input_path=source,
        vault_path=vault,
    )

    assert (result.created, result.updated, result.unchanged) == (1, 1, 1)
    assert result.written == result.created + result.updated + result.unchanged
