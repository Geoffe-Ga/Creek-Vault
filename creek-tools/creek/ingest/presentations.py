"""Presentation ingestor for the Creek pipeline.

Implements the PPTX half of §57 of the Creek Ontology — ingest
PowerPoint presentations as one fragment per file, with each slide
rendered as a labelled markdown section. Speaker notes are inlined
under their slide so reviewers can see what was said alongside the
visual content.

The presentation backend is decoupled through the
:class:`PresentationBackend` Protocol so callers can plug in any
implementation; tests inject a deterministic stub. The default
:class:`PythonPptxBackend` lazily imports ``python-pptx`` so the
rest of the package — and the unit tests — run cleanly even when
that optional dependency is not installed.

Optional dependency (install separately to enable PPTX support):

* ``python-pptx`` — pure-Python PPTX reader.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import TYPE_CHECKING, Any, Protocol, runtime_checkable

from creek.ingest.base import (
    Ingestor,
    ParsedFragment,
    RawDocument,
    file_modified_time,
    parse_authored_at,
)
from creek.models import SourcePlatform

if TYPE_CHECKING:
    from datetime import datetime
    from pathlib import Path

logger = logging.getLogger(__name__)


PRESENTATION_EXTENSIONS: frozenset[str] = frozenset({".pptx"})
"""File extensions handled by :class:`PresentationIngestor`."""


# ---- Public dataclasses + protocol --------------------------------------


@dataclass(frozen=True)
class SlideData:
    """A single slide's content extracted from a presentation.

    Attributes:
        index: 1-based slide number.
        title: Slide title text. ``None`` if no title placeholder.
        body: Concatenated body / placeholder text, newline-separated.
        notes: Speaker notes. Empty string when the slide has none.
    """

    index: int
    title: str | None
    body: str
    notes: str = ""


@dataclass(frozen=True)
class PresentationData:
    """A presentation decomposed into ordered slides.

    Attributes:
        title: Optional presentation title (from core properties).
        slides: Slides in presentation order.
    """

    title: str | None
    slides: tuple[SlideData, ...]


@runtime_checkable
class PresentationBackend(Protocol):
    """Pluggable presentation backend.

    Implementations must be deterministic and side-effect-free; the
    same input file should always produce the same
    :class:`PresentationData`.
    """

    def is_available(self) -> bool:
        """Return ``True`` when the backend can perform reads right now."""

    def read_presentation(self, path: Path) -> PresentationData:
        """Read *path* and return its decomposed :class:`PresentationData`."""


class PythonPptxUnavailableError(RuntimeError):
    """Raised when a PPTX read is attempted without ``python-pptx`` installed."""


# ---- Default backend ---------------------------------------------------


class PythonPptxBackend:
    """Presentation backend backed by ``python-pptx``.

    Imports of the optional dependency are deferred to call time so
    the rest of the package runs cleanly without it.
    """

    def is_available(self) -> bool:
        """Return ``True`` when ``python-pptx`` imports cleanly."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            return False
        return True

    def read_presentation(self, path: Path) -> PresentationData:
        """Read *path* and return a :class:`PresentationData`.

        Args:
            path: Filesystem path to a ``.pptx`` file.

        Returns:
            The decomposed :class:`PresentationData`.

        Raises:
            PythonPptxUnavailableError: When ``python-pptx`` is not
                installed.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            msg = (
                "python-pptx is required for PPTX ingestion. Install it with "
                "`pip install python-pptx`."
            )
            raise PythonPptxUnavailableError(msg) from exc

        prs = Presentation(str(path))
        slides: list[SlideData] = [
            _slide_to_data(slide, index)
            for index, slide in enumerate(prs.slides, start=1)
        ]
        title = self._extract_title(prs)
        return PresentationData(title=title, slides=tuple(slides))

    @staticmethod
    def _extract_title(prs: Any) -> str | None:
        """Return the presentation's core-properties title, if set.

        ``python-pptx`` exposes ``core_properties.title`` as either a
        non-empty string, an empty string, or ``None``. We coerce all
        absent-or-blank cases to ``None`` so callers can safely fall
        back to the file stem with ``data.title or path.stem`` —
        without the literal string ``"None"`` ever surfacing in
        rendered markdown or YAML frontmatter.
        """
        try:
            title = prs.core_properties.title
        except (AttributeError, ValueError):
            return None
        if title is None:
            return None
        stripped = str(title).strip()
        return stripped or None


def _extract_pptx_authored_at(path: Path) -> datetime | None:
    """Read the PPTX core-properties ``created`` (then ``modified``).

    FEAT-031: returns the deck's source-side authored date or
    ``None`` when neither field is populated, when ``python-pptx``
    is unavailable, or when the file is malformed. Filesystem mtime
    is never substituted — :attr:`Fragment.ingested` is the honest
    fallback chosen by downstream surfaces.
    """
    props = _open_pptx_core_properties(path)
    if props is None:
        return None
    for candidate in (
        getattr(props, "created", None),
        getattr(props, "modified", None),
    ):
        parsed = _safe_parse_authored_at(candidate)
        if parsed is not None:
            return parsed
    return None


def _open_pptx_core_properties(path: Path) -> Any | None:
    """Open *path* via python-pptx and return ``core_properties`` (or ``None``).

    Centralises the optional-import + file-open + attribute-lookup
    failure modes so :func:`_extract_pptx_authored_at` stays linear and
    below the per-function try-block cap.
    """
    try:  # noqa: TRY101  # Separate optional-dep absence from file-open failure.
        from pptx import Presentation as _Presentation
    except ImportError:
        return None
    try:  # noqa: TRY101  # Distinct failure mode from the ImportError branch.
        prs = _Presentation(str(path))
    except Exception:  # python-pptx raises a wide set on bad files
        logger.debug("Could not open PPTX core properties for %s", path)
        return None
    try:  # noqa: TRY101  # AttributeError is a distinct failure mode from file-open.
        return prs.core_properties
    except AttributeError:
        return None


def _safe_parse_authored_at(candidate: object) -> datetime | None:
    """Wrapper that swallows :class:`ValueError` from ``parse_authored_at``."""
    if candidate is None:
        return None
    try:
        return parse_authored_at(candidate)
    except ValueError:
        return None


def _slide_to_data(slide: Any, index: int) -> SlideData:
    """Convert a python-pptx slide object into a :class:`SlideData`."""
    title: str | None = None
    title_shape = slide.shapes.title
    body_chunks: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if title_shape is not None and shape is title_shape and title is None:
            title = text
        else:
            body_chunks.append(text)
    notes = ""
    if getattr(slide, "has_notes_slide", False):
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except AttributeError:
            notes = ""
    return SlideData(
        index=index,
        title=title,
        body="\n".join(body_chunks),
        notes=notes,
    )


# ---- PresentationIngestor ----------------------------------------------


class PresentationIngestor(Ingestor):
    """Ingest PPTX files as one fragment per presentation."""

    def __init__(self, backend: PresentationBackend | None = None) -> None:
        """Initialise with a backend; defaults to :class:`PythonPptxBackend`."""
        self.backend = backend if backend is not None else PythonPptxBackend()

    def discover(self, source_path: Path) -> list[RawDocument]:
        """Recursively find ``.pptx`` files under *source_path*.

        Single-file paths whose suffix is unsupported return an empty
        list. File bytes are not slurped at discovery time — the
        backend reads from disk via the :class:`Path`.
        """
        if source_path.is_file():
            paths = (
                [source_path]
                if source_path.suffix.lower() in PRESENTATION_EXTENSIONS
                else []
            )
        else:
            paths = [
                p
                for p in sorted(source_path.rglob("*"))
                if p.is_file() and p.suffix.lower() in PRESENTATION_EXTENSIONS
            ]
        return [
            RawDocument(
                path=path,
                content=b"",
                metadata={"original_file": str(path)},
                detected_encoding="binary",
            )
            for path in paths
        ]

    def parse(self, raw: RawDocument) -> list[ParsedFragment]:
        """Emit a single fragment carrying the entire presentation.

        FEAT-031: the PPTX core properties ``created`` (then
        ``modified``) become ``authored_at``. Falls through to
        ``None`` when the deck carries no creation date or when
        ``python-pptx`` is unavailable.
        """
        data = self.backend.read_presentation(raw.path)
        if not data.slides:
            logger.info(
                "Presentation %s has no slides; skipping fragment.",
                raw.path,
            )
            return []
        timestamp = file_modified_time(raw.path)
        authored_at = _extract_pptx_authored_at(raw.path)
        return [
            ParsedFragment(
                content="",  # Markdown rendered in convert_to_markdown.
                metadata={
                    "original_file": str(raw.path),
                    "title": data.title or raw.path.stem,
                    "slide_count": len(data.slides),
                    "authored_at": authored_at,
                },
                source_path=str(raw.path),
                timestamp=timestamp,
                payload=data,
            ),
        ]

    def convert_to_markdown(self, fragment: ParsedFragment) -> str:
        """Render the presentation as a markdown document.

        One ``## Slide N: Title`` block per slide with the body inlined
        and speaker notes (when present) under a ``**Speaker notes:**``
        sub-section. Reads the typed :class:`PresentationData` off
        :attr:`~creek.ingest.base.ParsedFragment.payload` rather than
        re-parsing a dict mirror in ``metadata`` (issue #166).
        """
        title = str(fragment.metadata.get("title", "Presentation"))
        data = fragment.payload
        slides: tuple[SlideData, ...] = (
            data.slides if isinstance(data, PresentationData) else ()
        )
        lines: list[str] = [f"# {title}", ""]
        for slide in slides:
            heading = f"## Slide {slide.index}"
            if slide.title:
                heading = f"{heading}: {slide.title}"
            lines.extend((heading, ""))
            if slide.body:
                lines.extend((slide.body, ""))
            if slide.notes:
                lines.extend(("**Speaker notes:**", "", slide.notes, ""))
        return "\n".join(lines).rstrip() + "\n"

    def generate_frontmatter(self, fragment: ParsedFragment) -> dict[str, Any]:
        """Produce YAML frontmatter for a presentation fragment.

        FEAT-031: the deck's PPTX core-properties ``created`` /
        ``modified`` land on ``authored_at`` as an ISO string when
        extracted; absent otherwise.
        """
        frontmatter_dict: dict[str, Any] = {
            "type": "fragment",
            "source": {
                "platform": SourcePlatform.PRESENTATION.value,
                "original_file": fragment.metadata.get(
                    "original_file",
                    fragment.source_path,
                ),
            },
            "title": fragment.metadata.get("title", ""),
            "slide_count": fragment.metadata.get("slide_count", 0),
            "ingested": fragment.timestamp.isoformat(),
        }
        authored_at: datetime | None = fragment.metadata.get("authored_at")
        if authored_at is not None:
            frontmatter_dict["authored_at"] = authored_at.isoformat()
        return frontmatter_dict


__all__ = [
    "PRESENTATION_EXTENSIONS",
    "PresentationBackend",
    "PresentationData",
    "PresentationIngestor",
    "PythonPptxBackend",
    "PythonPptxUnavailableError",
    "SlideData",
]
